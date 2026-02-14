# -*- coding: utf-8 -*-
"""Slides extras: Deep Learning com Sentinel-2"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# Paleta
AZUL_DARK = RGBColor(0x0d, 0x1b, 0x2a)
AZUL_MEDIO = RGBColor(0x1b, 0x26, 0x3b)
AZUL_ACCENT = RGBColor(0x41, 0x5a, 0x77)
CYAN = RGBColor(0x00, 0xd4, 0xff)
LARANJA = RGBColor(0xff, 0x6b, 0x35)
VERDE = RGBColor(0x00, 0xc9, 0xa7)
AMARELO = RGBColor(0xff, 0xc8, 0x00)
BRANCO = RGBColor(0xff, 0xff, 0xff)
CINZA_CLARO = RGBColor(0xe0, 0xe1, 0xdd)

def set_gradient(shape, color1, color2, angle=90):
    fill = shape.fill
    fill.gradient()
    fill.gradient_angle = angle
    fill.gradient_stops[0].color.rgb = color1
    fill.gradient_stops[1].color.rgb = color2

def add_content_slide(prs, title, bullets, accent_color=CYAN):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BRANCO
    bg.line.fill.background()

    sidebar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.2), prs.slide_height)
    set_gradient(sidebar, accent_color, AZUL_DARK, 180)
    sidebar.line.fill.background()

    # Decoracoes
    circ1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(12.5), Inches(-0.3), Inches(1.2), Inches(1.2))
    circ1.fill.solid()
    circ1.fill.fore_color.rgb = accent_color
    circ1.fill.fore_color.brightness = 0.85
    circ1.line.fill.background()

    circ2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(12.7), Inches(0.7), Inches(0.7), Inches(0.7))
    circ2.fill.solid()
    circ2.fill.fore_color.rgb = VERDE
    circ2.fill.fore_color.brightness = 0.85
    circ2.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(11), Inches(0.9))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = AZUL_DARK
    p.font.name = "Segoe UI"

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.15), Inches(1.5), Pt(4))
    line.fill.solid()
    line.fill.fore_color.rgb = accent_color
    line.line.fill.background()

    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(11.5), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if bullet.startswith("  "):
            p.text = "     " + bullet.strip()
            p.font.size = Pt(17)
            p.font.color.rgb = AZUL_ACCENT
        elif bullet == "":
            p.text = ""
            p.font.size = Pt(10)
        else:
            p.text = bullet
            p.font.size = Pt(19)
            p.font.color.rgb = AZUL_DARK
        p.font.name = "Segoe UI"
        p.space_after = Pt(8)

    return slide

def add_comparison_slide(prs, title, left_title, left_items, right_title, right_items, left_color=CYAN, right_color=LARANJA):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BRANCO
    bg.line.fill.background()

    sidebar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.2), prs.slide_height)
    set_gradient(sidebar, VERDE, CYAN, 180)
    sidebar.line.fill.background()

    # Divisor central
    divider = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.4), Inches(1.5), Inches(0.1), Inches(5.5))
    set_gradient(divider, left_color, right_color, 180)
    divider.line.fill.background()

    # Titulo principal
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = AZUL_DARK
    p.font.name = "Segoe UI"

    # Coluna esquerda
    icon_left = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), Inches(1.3), Inches(0.5), Inches(0.5))
    icon_left.fill.solid()
    icon_left.fill.fore_color.rgb = left_color
    icon_left.line.fill.background()

    left_title_box = slide.shapes.add_textbox(Inches(1.1), Inches(1.35), Inches(5), Inches(0.5))
    tf = left_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = left_color
    p.font.name = "Segoe UI"

    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(5.7), Inches(5))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(15)
        p.font.color.rgb = AZUL_DARK
        p.font.name = "Segoe UI"
        p.space_after = Pt(6)

    # Coluna direita
    icon_right = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.7), Inches(1.3), Inches(0.5), Inches(0.5))
    icon_right.fill.solid()
    icon_right.fill.fore_color.rgb = right_color
    icon_right.line.fill.background()

    right_title_box = slide.shapes.add_textbox(Inches(7.3), Inches(1.35), Inches(5), Inches(0.5))
    tf = right_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = right_color
    p.font.name = "Segoe UI"

    right_box = slide.shapes.add_textbox(Inches(6.7), Inches(2), Inches(5.7), Inches(5))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(right_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(15)
        p.font.color.rgb = AZUL_DARK
        p.font.name = "Segoe UI"
        p.space_after = Pt(6)

    return slide

def add_table_slide(prs, title, headers, rows, accent_color=CYAN):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BRANCO
    bg.line.fill.background()

    sidebar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.2), prs.slide_height)
    set_gradient(sidebar, accent_color, AZUL_DARK, 180)
    sidebar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(11), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = AZUL_DARK
    p.font.name = "Segoe UI"

    num_cols = len(headers)
    num_rows = len(rows) + 1
    table = slide.shapes.add_table(num_rows, num_cols, Inches(0.5), Inches(1.4), Inches(12), Inches(0.55 * num_rows)).table

    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = AZUL_DARK
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = BRANCO
        p.alignment = PP_ALIGN.CENTER
        p.font.name = "Segoe UI"

    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = str(val)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(13)
            p.font.color.rgb = AZUL_DARK
            p.alignment = PP_ALIGN.CENTER
            p.font.name = "Segoe UI"
            if ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xf0, 0xf8, 0xff)

    return slide

def add_diagram_slide(prs):
    """Slide com diagrama visual: Indices vs CNN"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BRANCO
    bg.line.fill.background()

    sidebar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.2), prs.slide_height)
    set_gradient(sidebar, LARANJA, AZUL_DARK, 180)
    sidebar.line.fill.background()

    # Titulo
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Comparacao Visual: Indices vs Imagem Completa"
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = AZUL_DARK
    p.font.name = "Segoe UI"

    # === ABORDAGEM ATUAL (LightGBM) ===
    # Box superior
    box1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(1.1), Inches(12.5), Inches(2.7))
    box1.fill.solid()
    box1.fill.fore_color.rgb = RGBColor(0xe8, 0xf4, 0xf8)
    box1.line.color.rgb = CYAN
    box1.line.width = Pt(2)

    # Label
    label1 = slide.shapes.add_textbox(Inches(0.6), Inches(1.2), Inches(4), Inches(0.5))
    tf = label1.text_frame
    p = tf.paragraphs[0]
    p.text = "ABORDAGEM ATUAL (LightGBM)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = CYAN
    p.font.name = "Segoe UI"

    # Sentinel icon
    sent1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.9), Inches(1.8), Inches(1.5))
    sent1.fill.solid()
    sent1.fill.fore_color.rgb = AZUL_MEDIO
    sent1.line.fill.background()

    sent1_txt = slide.shapes.add_textbox(Inches(0.8), Inches(2.3), Inches(1.8), Inches(0.8))
    tf = sent1_txt.text_frame
    p = tf.paragraphs[0]
    p.text = "Sentinel-2\n13 bandas"
    p.font.size = Pt(12)
    p.font.color.rgb = BRANCO
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Segoe UI"

    # Arrow 1
    arrow1 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(2.8), Inches(2.4), Inches(0.8), Inches(0.4))
    arrow1.fill.solid()
    arrow1.fill.fore_color.rgb = AZUL_ACCENT
    arrow1.line.fill.background()

    # Calculo indices
    calc1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.8), Inches(1.9), Inches(2.2), Inches(1.5))
    calc1.fill.solid()
    calc1.fill.fore_color.rgb = VERDE
    calc1.line.fill.background()

    calc1_txt = slide.shapes.add_textbox(Inches(3.8), Inches(2.1), Inches(2.2), Inches(1.2))
    tf = calc1_txt.text_frame
    p = tf.paragraphs[0]
    p.text = "Calcula\nNDVI, NBR\ndNBR"
    p.font.size = Pt(12)
    p.font.color.rgb = BRANCO
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Segoe UI"

    # Arrow 2
    arrow2 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.2), Inches(2.4), Inches(0.8), Inches(0.4))
    arrow2.fill.solid()
    arrow2.fill.fore_color.rgb = AZUL_ACCENT
    arrow2.line.fill.background()

    # Features
    feat1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.2), Inches(1.9), Inches(2.2), Inches(1.5))
    feat1.fill.solid()
    feat1.fill.fore_color.rgb = LARANJA
    feat1.line.fill.background()

    feat1_txt = slide.shapes.add_textbox(Inches(7.2), Inches(2.1), Inches(2.2), Inches(1.2))
    tf = feat1_txt.text_frame
    p = tf.paragraphs[0]
    p.text = "~20 numeros\npor hotspot"
    p.font.size = Pt(12)
    p.font.color.rgb = BRANCO
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Segoe UI"

    # Arrow 3
    arrow3 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(9.6), Inches(2.4), Inches(0.8), Inches(0.4))
    arrow3.fill.solid()
    arrow3.fill.fore_color.rgb = AZUL_ACCENT
    arrow3.line.fill.background()

    # LightGBM
    lgbm = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.6), Inches(1.9), Inches(2), Inches(1.5))
    lgbm.fill.solid()
    lgbm.fill.fore_color.rgb = AZUL_DARK
    lgbm.line.fill.background()

    lgbm_txt = slide.shapes.add_textbox(Inches(10.6), Inches(2.3), Inches(2), Inches(0.8))
    tf = lgbm_txt.text_frame
    p = tf.paragraphs[0]
    p.text = "LightGBM\n82%"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = BRANCO
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Segoe UI"

    # === ABORDAGEM DEEP LEARNING ===
    # Box inferior
    box2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(4.0), Inches(12.5), Inches(2.7))
    box2.fill.solid()
    box2.fill.fore_color.rgb = RGBColor(0xff, 0xf0, 0xe8)
    box2.line.color.rgb = LARANJA
    box2.line.width = Pt(2)

    # Label
    label2 = slide.shapes.add_textbox(Inches(0.6), Inches(4.1), Inches(4), Inches(0.5))
    tf = label2.text_frame
    p = tf.paragraphs[0]
    p.text = "ABORDAGEM FUTURA (Deep Learning)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = LARANJA
    p.font.name = "Segoe UI"

    # Sentinel icon 2
    sent2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(4.8), Inches(1.8), Inches(1.5))
    sent2.fill.solid()
    sent2.fill.fore_color.rgb = AZUL_MEDIO
    sent2.line.fill.background()

    sent2_txt = slide.shapes.add_textbox(Inches(0.8), Inches(5.2), Inches(1.8), Inches(0.8))
    tf = sent2_txt.text_frame
    p = tf.paragraphs[0]
    p.text = "Sentinel-2\n13 bandas"
    p.font.size = Pt(12)
    p.font.color.rgb = BRANCO
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Segoe UI"

    # Arrow 1
    arrow4 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(2.8), Inches(5.3), Inches(0.8), Inches(0.4))
    arrow4.fill.solid()
    arrow4.fill.fore_color.rgb = AZUL_ACCENT
    arrow4.line.fill.background()

    # Patch
    patch = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.8), Inches(4.8), Inches(2.2), Inches(1.5))
    patch.fill.solid()
    patch.fill.fore_color.rgb = AMARELO
    patch.line.fill.background()

    patch_txt = slide.shapes.add_textbox(Inches(3.8), Inches(5.0), Inches(2.2), Inches(1.2))
    tf = patch_txt.text_frame
    p = tf.paragraphs[0]
    p.text = "Recorta Patch\n64x64 pixels"
    p.font.size = Pt(12)
    p.font.color.rgb = AZUL_DARK
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Segoe UI"

    # Arrow 2
    arrow5 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.2), Inches(5.3), Inches(0.8), Inches(0.4))
    arrow5.fill.solid()
    arrow5.fill.fore_color.rgb = AZUL_ACCENT
    arrow5.line.fill.background()

    # Tensor
    tensor = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.2), Inches(4.8), Inches(2.2), Inches(1.5))
    tensor.fill.solid()
    tensor.fill.fore_color.rgb = LARANJA
    tensor.line.fill.background()

    tensor_txt = slide.shapes.add_textbox(Inches(7.2), Inches(5.0), Inches(2.2), Inches(1.2))
    tf = tensor_txt.text_frame
    p = tf.paragraphs[0]
    p.text = "~53.000 valores\n(64x64x13)"
    p.font.size = Pt(12)
    p.font.color.rgb = BRANCO
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Segoe UI"

    # Arrow 3
    arrow6 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(9.6), Inches(5.3), Inches(0.8), Inches(0.4))
    arrow6.fill.solid()
    arrow6.fill.fore_color.rgb = AZUL_ACCENT
    arrow6.line.fill.background()

    # CNN
    cnn = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.6), Inches(4.8), Inches(2), Inches(1.5))
    cnn.fill.solid()
    cnn.fill.fore_color.rgb = VERDE
    cnn.line.fill.background()

    cnn_txt = slide.shapes.add_textbox(Inches(10.6), Inches(5.2), Inches(2), Inches(0.8))
    tf = cnn_txt.text_frame
    p = tf.paragraphs[0]
    p.text = "CNN\n> 82%?"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = BRANCO
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Segoe UI"

    return slide

def add_example_slide(prs):
    """Slide com exemplo pratico de dois hotspots"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BRANCO
    bg.line.fill.background()

    sidebar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.2), prs.slide_height)
    set_gradient(sidebar, VERDE, AZUL_DARK, 180)
    sidebar.line.fill.background()

    # Titulo
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Exemplo: Por que CNN pode ser melhor?"
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = AZUL_DARK
    p.font.name = "Segoe UI"

    # Subtitulo
    sub = slide.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(12), Inches(0.5))
    tf = sub.text_frame
    p = tf.paragraphs[0]
    p.text = "Dois hotspots com MESMO NDVI = 0.45, mas contextos diferentes"
    p.font.size = Pt(16)
    p.font.color.rgb = AZUL_ACCENT
    p.font.name = "Segoe UI"

    # === HOTSPOT A ===
    box_a = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.6), Inches(5.8), Inches(4))
    box_a.fill.solid()
    box_a.fill.fore_color.rgb = RGBColor(0xd6, 0xf5, 0xd6)
    box_a.line.color.rgb = VERDE
    box_a.line.width = Pt(2)

    # Titulo A
    title_a = slide.shapes.add_textbox(Inches(0.7), Inches(1.8), Inches(5.4), Inches(0.5))
    tf = title_a.text_frame
    p = tf.paragraphs[0]
    p.text = "HOTSPOT A"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = VERDE
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Segoe UI"

    # Grid A (area queimada grande)
    grid_a = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(2.5), Inches(3.5), Inches(1.8))
    grid_a.fill.solid()
    grid_a.fill.fore_color.rgb = RGBColor(0x8b, 0x45, 0x13)  # Marrom (queimado)
    grid_a.line.color.rgb = AZUL_DARK
    grid_a.line.width = Pt(1)

    # Label A
    label_a = slide.shapes.add_textbox(Inches(0.7), Inches(4.5), Inches(5.4), Inches(0.9))
    tf = label_a.text_frame
    p = tf.paragraphs[0]
    p.text = "Area queimada GRANDE e uniforme"
    p.font.size = Pt(14)
    p.font.color.rgb = AZUL_DARK
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Segoe UI"
    p = tf.add_paragraph()
    p.text = "Provavelmente REAL"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = VERDE
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Segoe UI"

    # === HOTSPOT B ===
    box_b = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(1.6), Inches(5.8), Inches(4))
    box_b.fill.solid()
    box_b.fill.fore_color.rgb = RGBColor(0xf5, 0xd6, 0xd6)
    box_b.line.color.rgb = LARANJA
    box_b.line.width = Pt(2)

    # Titulo B
    title_b = slide.shapes.add_textbox(Inches(7.2), Inches(1.8), Inches(5.4), Inches(0.5))
    tf = title_b.text_frame
    p = tf.paragraphs[0]
    p.text = "HOTSPOT B"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = LARANJA
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Segoe UI"

    # Grid B (fundo verde com ponto pequeno)
    grid_b_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8), Inches(2.5), Inches(3.5), Inches(1.8))
    grid_b_bg.fill.solid()
    grid_b_bg.fill.fore_color.rgb = RGBColor(0x22, 0x8b, 0x22)  # Verde (vegetacao)
    grid_b_bg.line.color.rgb = AZUL_DARK
    grid_b_bg.line.width = Pt(1)

    # Ponto isolado
    point_b = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.5), Inches(3.1), Inches(0.5), Inches(0.5))
    point_b.fill.solid()
    point_b.fill.fore_color.rgb = RGBColor(0x8b, 0x45, 0x13)  # Marrom
    point_b.line.fill.background()

    # Label B
    label_b = slide.shapes.add_textbox(Inches(7.2), Inches(4.5), Inches(5.4), Inches(0.9))
    tf = label_b.text_frame
    p = tf.paragraphs[0]
    p.text = "Ponto ISOLADO em area verde"
    p.font.size = Pt(14)
    p.font.color.rgb = AZUL_DARK
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Segoe UI"
    p = tf.add_paragraph()
    p.text = "Provavelmente ESPURIO"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = LARANJA
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Segoe UI"

    # Conclusao
    concl = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.8), Inches(12.3), Inches(1.3))
    concl.fill.solid()
    concl.fill.fore_color.rgb = AZUL_DARK
    concl.line.fill.background()

    concl_txt = slide.shapes.add_textbox(Inches(0.5), Inches(6.0), Inches(12.3), Inches(1))
    tf = concl_txt.text_frame
    p = tf.paragraphs[0]
    p.text = "LightGBM: Ve apenas NDVI=0.45 -> Mesma classificacao para ambos"
    p.font.size = Pt(15)
    p.font.color.rgb = CINZA_CLARO
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Segoe UI"
    p = tf.add_paragraph()
    p.text = "CNN: Ve o PADRAO ESPACIAL -> Consegue distinguir os dois!"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = CYAN
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Segoe UI"

    return slide

def create_slides():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # SLIDE 1: Introducao Deep Learning
    add_content_slide(prs, "Trabalho Futuro: Deep Learning com Sentinel-2", [
        "Abordagem atual (LightGBM):",
        "  Extrai INDICES espectrais das imagens Sentinel-2",
        "  NDVI, NBR, dNBR -> ~20 numeros por hotspot",
        "  Perde informacao de CONTEXTO ESPACIAL",
        "",
        "Proposta futura (CNN - Redes Neurais Convolucionais):",
        "  Usar a IMAGEM COMPLETA ao redor do hotspot",
        "  Patch de 64x64 pixels com 13 bandas",
        "  Preserva padroes de textura, forma e vizinhanca",
        "  Potencial para superar 82% de acuracia"
    ], LARANJA)

    # SLIDE 2: O que sao indices espectrais
    add_content_slide(prs, "O que sao Indices Espectrais?", [
        "Formulas matematicas que combinam bandas do satelite:",
        "",
        "NDVI = (NIR - RED) / (NIR + RED)",
        "  Mede saude da vegetacao (-1 a +1)",
        "  Vegetacao saudavel: NDVI alto (~0.6-0.9)",
        "  Solo exposto/queimado: NDVI baixo (~0.1-0.3)",
        "",
        "NBR = (NIR - SWIR) / (NIR + SWIR)",
        "  Mede severidade de queimada",
        "",
        "dNBR = NBR_antes - NBR_depois",
        "  Detecta MUDANCA apos fogo"
    ], CYAN)

    # SLIDE 3: Diagrama comparativo
    add_diagram_slide(prs)

    # SLIDE 4: O que se perde
    add_table_slide(prs, "Informacao Perdida ao Usar Apenas Indices",
        ["Informacao", "Com NDVI/NBR", "Com Imagem CNN"],
        [
            ["Contexto espacial", "PERDIDO", "PRESERVADO"],
            ["Padroes de textura", "PERDIDO", "PRESERVADO"],
            ["Forma da queimada", "PERDIDO", "PRESERVADO"],
            ["Vizinhanca do hotspot", "PERDIDO", "PRESERVADO"],
            ["Bordas e contornos", "PERDIDO", "PRESERVADO"]
        ], LARANJA)

    # SLIDE 5: Exemplo pratico
    add_example_slide(prs)

    # SLIDE 6: Comparacao final
    add_table_slide(prs, "Comparacao: LightGBM vs CNN",
        ["Aspecto", "Indices (LightGBM)", "Imagem (CNN)"],
        [
            ["Dados por hotspot", "~20 numeros", "~53.000 valores"],
            ["Contexto espacial", "Nao", "Sim"],
            ["Velocidade", "Rapido (CPU)", "Lento (GPU)"],
            ["Hardware", "CPU ok", "GPU necessaria"],
            ["Interpretabilidade", "Alta", "Baixa"],
            ["Acuracia atual", "82%", "A ser testado"]
        ], VERDE)

    # SLIDE 7: Proximos passos
    add_content_slide(prs, "Proximos Passos para Deep Learning", [
        "1. Preparacao de dados:",
        "  Recortar patches 64x64 centrados em cada hotspot",
        "  Criar dataset de imagens rotuladas",
        "",
        "2. Arquiteturas a testar:",
        "  CNN simples (Conv2D + MaxPool + Dense)",
        "  ResNet pre-treinada (transfer learning)",
        "  U-Net para segmentacao de areas queimadas",
        "",
        "3. Requisitos:",
        "  GPU (Google Colab, AWS, ou local)",
        "  Mais dados de treinamento",
        "  Tempo de experimentacao"
    ], VERDE)

    prs.save("Slides_Deep_Learning_Sentinel2.pptx")
    print("Slides salvos: Slides_Deep_Learning_Sentinel2.pptx")
    print(f"Total: {len(prs.slides)} slides")

if __name__ == "__main__":
    print("Criando slides sobre Deep Learning...")
    create_slides()
    print("Concluido!")
