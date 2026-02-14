# -*- coding: utf-8 -*-
"""Apresentacao moderna para qualificacao de mestrado"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn
from docx import Document
from docx.shared import Pt as DocxPt
from docx.enum.text import WD_ALIGN_PARAGRAPH

# Paleta moderna
AZUL_DARK = RGBColor(0x0d, 0x1b, 0x2a)      # #0d1b2a - quase preto azulado
AZUL_MEDIO = RGBColor(0x1b, 0x26, 0x3b)     # #1b263b
AZUL_ACCENT = RGBColor(0x41, 0x5a, 0x77)    # #415a77
AZUL_LIGHT = RGBColor(0x77, 0x8d, 0xa9)     # #778da9
CYAN = RGBColor(0x00, 0xd4, 0xff)           # #00d4ff - destaque cyan
LARANJA = RGBColor(0xff, 0x6b, 0x35)        # #ff6b35 - destaque laranja
VERDE = RGBColor(0x00, 0xc9, 0xa7)          # #00c9a7 - verde agua
BRANCO = RGBColor(0xff, 0xff, 0xff)
CINZA_CLARO = RGBColor(0xe0, 0xe1, 0xdd)    # #e0e1dd

def set_shape_gradient(shape, color1, color2, angle=90):
    """Aplica gradiente a uma forma"""
    fill = shape.fill
    fill.gradient()
    fill.gradient_angle = angle
    fill.gradient_stops[0].color.rgb = color1
    fill.gradient_stops[1].color.rgb = color2

def add_modern_title_slide(prs, title, subtitle, author, advisor, coadvisor):
    """Slide de titulo moderno com gradiente"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Fundo gradiente
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    set_shape_gradient(bg, AZUL_DARK, AZUL_MEDIO, 135)
    bg.line.fill.background()

    # Forma geometrica decorativa (triangulo)
    tri1 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Inches(10), Inches(0), Inches(3.5), Inches(4))
    tri1.fill.solid()
    tri1.fill.fore_color.rgb = RGBColor(0x1b, 0x26, 0x3b)
    tri1.fill.fore_color.brightness = 0.1
    tri1.line.fill.background()

    # Linha cyan decorativa
    line1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.8), Inches(1.5), Pt(4))
    line1.fill.solid()
    line1.fill.fore_color.rgb = CYAN
    line1.line.fill.background()

    # Titulo principal
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = BRANCO
    p.font.name = "Segoe UI"

    # Subtitulo
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(0.8))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(18)
    p.font.color.rgb = CYAN
    p.font.name = "Segoe UI Light"

    # Info do autor
    info_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(8), Inches(1.5))
    tf = info_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"Mestrando: {author}"
    p.font.size = Pt(14)
    p.font.color.rgb = CINZA_CLARO
    p.font.name = "Segoe UI"

    p = tf.add_paragraph()
    p.text = f"Orientador: {advisor}"
    p.font.size = Pt(14)
    p.font.color.rgb = CINZA_CLARO
    p.font.name = "Segoe UI"

    p = tf.add_paragraph()
    p.text = f"Coorientador: {coadvisor}"
    p.font.size = Pt(14)
    p.font.color.rgb = CINZA_CLARO
    p.font.name = "Segoe UI"

    return slide

def add_section_slide(prs, number, title):
    """Slide de secao com numero grande"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Fundo
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    set_shape_gradient(bg, AZUL_DARK, AZUL_MEDIO, 45)
    bg.line.fill.background()

    # Circulo com numero
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), Inches(2), Inches(2.5), Inches(2.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = CYAN
    circle.line.fill.background()

    # Numero
    num_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(2.5), Inches(2))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{number:02d}"
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = AZUL_DARK
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Segoe UI Black"

    # Titulo da secao
    title_box = slide.shapes.add_textbox(Inches(4), Inches(2.5), Inches(8.5), Inches(2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = BRANCO
    p.font.name = "Segoe UI"

    return slide

def add_content_slide(prs, title, bullets, highlight_color=CYAN):
    """Slide de conteudo com bullets modernos"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Fundo branco
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BRANCO
    bg.line.fill.background()

    # Barra lateral esquerda
    sidebar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.15), prs.slide_height)
    sidebar.fill.solid()
    sidebar.fill.fore_color.rgb = highlight_color
    sidebar.line.fill.background()

    # Titulo
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12), Inches(0.9))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = AZUL_DARK
    p.font.name = "Segoe UI"

    # Linha sob titulo
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.25), Inches(2), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = highlight_color
    line.line.fill.background()

    # Bullets
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(12.3), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        # Verifica se eh sub-item
        if bullet.startswith("  "):
            p.text = "    " + bullet.strip()
            p.font.size = Pt(18)
            p.font.color.rgb = AZUL_ACCENT
        else:
            p.text = bullet
            p.font.size = Pt(20)
            p.font.color.rgb = AZUL_DARK

        p.font.name = "Segoe UI"
        p.space_after = Pt(10)

    return slide

def add_metrics_slide(prs, title, metrics):
    """Slide com cards de metricas"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Fundo
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    set_shape_gradient(bg, AZUL_DARK, AZUL_MEDIO, 180)
    bg.line.fill.background()

    # Titulo
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12), Inches(0.9))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = BRANCO
    p.font.name = "Segoe UI"

    # Cards
    card_width = Inches(2.8)
    card_height = Inches(3)
    start_x = Inches(0.7)
    y_pos = Inches(2.2)
    spacing = Inches(3.1)

    colors = [CYAN, VERDE, LARANJA, AZUL_LIGHT]

    for i, (name, value) in enumerate(metrics):
        x = start_x + (i * spacing)

        # Card com cantos arredondados
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y_pos, card_width, card_height)
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(0x1b, 0x26, 0x3b)
        card.line.color.rgb = colors[i % 4]
        card.line.width = Pt(3)

        # Barra colorida no topo do card
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y_pos, card_width, Inches(0.1))
        bar.fill.solid()
        bar.fill.fore_color.rgb = colors[i % 4]
        bar.line.fill.background()

        # Valor
        val_box = slide.shapes.add_textbox(x, y_pos + Inches(0.5), card_width, Inches(1.2))
        tf = val_box.text_frame
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = colors[i % 4]
        p.alignment = PP_ALIGN.CENTER
        p.font.name = "Segoe UI Black"

        # Nome
        name_box = slide.shapes.add_textbox(x, y_pos + Inches(1.8), card_width, Inches(1))
        tf = name_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(16)
        p.font.color.rgb = CINZA_CLARO
        p.alignment = PP_ALIGN.CENTER
        p.font.name = "Segoe UI"

    return slide

def add_table_slide(prs, title, headers, rows):
    """Slide com tabela moderna"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Fundo
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BRANCO
    bg.line.fill.background()

    # Barra lateral
    sidebar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.15), prs.slide_height)
    sidebar.fill.solid()
    sidebar.fill.fore_color.rgb = CYAN
    sidebar.line.fill.background()

    # Titulo
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = AZUL_DARK
    p.font.name = "Segoe UI"

    # Tabela
    num_cols = len(headers)
    num_rows = len(rows) + 1

    x = Inches(0.5)
    y = Inches(1.5)
    cx = Inches(12.3)
    cy = Inches(0.55 * num_rows)

    table = slide.shapes.add_table(num_rows, num_cols, x, y, cx, cy).table

    # Cabecalho
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = AZUL_DARK
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = BRANCO
        p.alignment = PP_ALIGN.CENTER
        p.font.name = "Segoe UI"

    # Dados
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = str(val)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(13)
            p.font.color.rgb = AZUL_DARK
            p.alignment = PP_ALIGN.CENTER
            p.font.name = "Segoe UI"

            if ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xf0, 0xf4, 0xf8)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = BRANCO

    return slide

def add_two_column_slide(prs, title, left_title, left_items, right_title, right_items):
    """Slide com duas colunas"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Fundo
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BRANCO
    bg.line.fill.background()

    # Barra lateral
    sidebar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.15), prs.slide_height)
    sidebar.fill.solid()
    sidebar.fill.fore_color.rgb = VERDE
    sidebar.line.fill.background()

    # Titulo principal
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = AZUL_DARK
    p.font.name = "Segoe UI"

    # Coluna esquerda - titulo
    left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(6), Inches(0.6))
    tf = left_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = CYAN
    p.font.name = "Segoe UI"

    # Coluna esquerda - items
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(6), Inches(5))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = AZUL_DARK
        p.font.name = "Segoe UI"
        p.space_after = Pt(8)

    # Coluna direita - titulo
    right_title_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.2), Inches(6), Inches(0.6))
    tf = right_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = LARANJA
    p.font.name = "Segoe UI"

    # Coluna direita - items
    right_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(6), Inches(5))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(right_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = AZUL_DARK
        p.font.name = "Segoe UI"
        p.space_after = Pt(8)

    return slide

def add_end_slide(prs):
    """Slide final"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Fundo
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    set_shape_gradient(bg, AZUL_DARK, AZUL_MEDIO, 135)
    bg.line.fill.background()

    # Texto principal
    title_box = slide.shapes.add_textbox(Inches(0), Inches(2.5), prs.slide_width, Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Obrigado!"
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = BRANCO
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Segoe UI"

    # Subtexto
    sub_box = slide.shapes.add_textbox(Inches(0), Inches(4.2), prs.slide_width, Inches(0.8))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Perguntas?"
    p.font.size = Pt(28)
    p.font.color.rgb = CYAN
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Segoe UI Light"

    # Contato
    contact_box = slide.shapes.add_textbox(Inches(0), Inches(5.5), prs.slide_width, Inches(1))
    tf = contact_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Andrey Ruben Ribeiro Bessa"
    p.font.size = Pt(16)
    p.font.color.rgb = CINZA_CLARO
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Segoe UI"

    p = tf.add_paragraph()
    p.text = "Programa de Pos-Graduacao em Engenharia Eletrica"
    p.font.size = Pt(14)
    p.font.color.rgb = AZUL_LIGHT
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Segoe UI"

    return slide


def create_presentation():
    """Cria a apresentacao completa"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Dados
    author = "Andrey Ruben Ribeiro Bessa"
    advisor = "Prof. D.Sc. Celso Barbosa Carvalho"
    coadvisor = "Prof. D.Sc. Andre Chaves Mendes"

    # SLIDE 1: Titulo
    add_modern_title_slide(prs,
        "Classificacao Automatica de Deteccoes\nde Focos de Incendio na Regiao\nMATOPIBA Usando Machine Learning",
        "Qualificacao de Mestrado | Engenharia Eletrica",
        author, advisor, coadvisor)

    # SLIDE 2: Agenda
    add_content_slide(prs, "Agenda", [
        "01  Introducao e Motivacao",
        "02  Objetivos do Trabalho",
        "03  Fundamentacao Teorica",
        "04  Metodologia Proposta",
        "05  Resultados Experimentais",
        "06  Cronograma",
        "07  Conclusoes e Proximos Passos"
    ])

    # === SECAO 1: INTRODUCAO ===
    add_section_slide(prs, 1, "Introducao e\nMotivacao")

    add_content_slide(prs, "Contexto: Incendios Florestais no Brasil", [
        "Brasil: um dos paises com maior incidencia de incendios florestais",
        "INPE monitora milhares de focos de calor diariamente via satelite",
        "Sistema FIRMS/NASA gera alertas automaticos em tempo quase-real",
        "  Deteccoes baseadas em anomalias termicas (MODIS, VIIRS)",
        "  Cobertura global com atualizacao a cada 3 horas",
        "Problema critico: Alta taxa de FALSOS POSITIVOS",
        "  Sobrecarrega orgaos de monitoramento ambiental",
        "  Verificacao manual inviavel em escala operacional"
    ], CYAN)

    add_content_slide(prs, "O Problema: Deteccoes Espurias", [
        "Fontes comuns de falsos positivos:",
        "  Reflexos solares em superficies metalicas (telhados, silos)",
        "  Queima controlada de gas em industrias",
        "  Alvos urbanos com alta temperatura aparente",
        "  Nuvens quentes e bordas de nuvens",
        "  Solo exposto com alta reflectancia",
        "",
        "Impacto operacional:",
        "  Desperdicio de recursos humanos e financeiros",
        "  Atrasos na resposta a incendios reais",
        "  Perda de credibilidade do sistema de alertas"
    ], LARANJA)

    add_content_slide(prs, "Area de Estudo: Regiao MATOPIBA", [
        "Fronteira agricola estrategica no Cerrado brasileiro",
        "Abrange partes de 4 estados:",
        "  Maranhao | Tocantins | Piaui | Bahia",
        "",
        "Caracteristicas relevantes:",
        "  Alta incidencia de queimadas (naturais e agricolas)",
        "  Expansao agropecuaria intensa",
        "  Relevancia economica e ambiental significativa",
        "  Disponibilidade de dados satelitais (2022-2024)"
    ], VERDE)

    # === SECAO 2: OBJETIVOS ===
    add_section_slide(prs, 2, "Objetivos")

    add_content_slide(prs, "Objetivo Geral", [
        "Desenvolver um sistema de classificacao automatica",
        "baseado em aprendizado de maquina para:",
        "",
        "  Distinguir deteccoes de incendios REAIS de ESPURIAS",
        "  nos dados do sistema FIRMS/NASA",
        "",
        "Aplicacao: Regiao MATOPIBA",
        "Periodo: 2022-2024",
        "Algoritmo: Gradient Boosting (LightGBM)"
    ], CYAN)

    add_content_slide(prs, "Objetivos Especificos", [
        "1. Implementar pipeline de ingestao de dados geoespaciais multifonte",
        "2. Desenvolver metodologia de weak labeling automatica",
        "  Usar produto MCD64A1 como referencia",
        "3. Definir e extrair caracteristicas relevantes para classificacao",
        "  Features temporais, meteorologicas, de vegetacao e espaciais",
        "4. Treinar e otimizar modelos de gradient boosting",
        "5. Validar com estrategias espaciais e temporais rigorosas",
        "6. Analisar importancia das caracteristicas preditivas"
    ], CYAN)

    # === SECAO 3: FUNDAMENTACAO ===
    add_section_slide(prs, 3, "Fundamentacao\nTeorica")

    add_content_slide(prs, "Sensoriamento Remoto: Deteccao de Incendios", [
        "Principio: Deteccao de anomalias termicas via infravermelho",
        "",
        "Satelites utilizados:",
        "  MODIS (Terra/Aqua) - resolucao 1km, 4x ao dia",
        "  VIIRS (Suomi NPP/NOAA-20) - resolucao 375m, 2x ao dia",
        "",
        "Produtos de dados:",
        "  FIRMS: Focos ativos (pontos de calor) em tempo quase-real",
        "  MCD64A1: Produto de area queimada mensal (500m)",
        "  Sentinel-2: Imagens multiespectrais de alta resolucao (10-20m)"
    ], VERDE)

    add_content_slide(prs, "Machine Learning: Gradient Boosting", [
        "Tecnica de ensemble baseada em arvores de decisao",
        "",
        "Funcionamento:",
        "  Arvores treinadas sequencialmente",
        "  Cada arvore corrige erros das anteriores",
        "  Combinacao ponderada das predicoes",
        "",
        "Algoritmos utilizados:",
        "  LightGBM: Otimizado para velocidade e memoria",
        "  XGBoost: Referencia em competicoes de ML",
        "",
        "Vantagens: Alta performance, interpretabilidade, robustez"
    ], VERDE)

    add_content_slide(prs, "Weak Labeling: Supervisao Fraca", [
        "Problema: Rotulacao manual e inviavel em larga escala",
        "",
        "Solucao: Gerar rotulos automaticamente usando dados auxiliares",
        "",
        "Abordagem proposta:",
        "  Comparar deteccoes FIRMS com produto MCD64A1 (area queimada)",
        "  Janela temporal: +/- 15 dias da deteccao",
        "  Raio espacial: 1 km ao redor do ponto",
        "",
        "Classificacao:",
        "  REAL (1): Correspondencia encontrada no MCD64A1",
        "  ESPURIA (0): Sem correspondencia espaciotemporal"
    ], VERDE)

    # === SECAO 4: METODOLOGIA ===
    add_section_slide(prs, 4, "Metodologia")

    add_content_slide(prs, "Arquitetura do Pipeline", [
        "Sistema modular em 4 etapas:",
        "",
        "1. INGESTAO",
        "  Download automatizado de dados brutos via APIs",
        "2. PROCESSAMENTO",
        "  Padronizacao, reprojecao e alinhamento espaciotemporal",
        "3. FEATURES",
        "  Extracao de caracteristicas + weak labeling automatico",
        "4. MODELAGEM",
        "  Treinamento, validacao cruzada e inferencia"
    ], LARANJA)

    add_table_slide(prs, "Fontes de Dados Integradas",
        ["Fonte", "Tipo", "Resolucao", "Frequencia"],
        [
            ["FIRMS/VIIRS", "Focos de calor", "375m", "~3 horas"],
            ["MCD64A1", "Area queimada", "500m", "Mensal"],
            ["Sentinel-2", "Multiespectral", "10-20m", "5 dias"],
            ["ERA5-Land", "Meteorologia", "~9km", "Horario"]
        ])

    add_table_slide(prs, "Caracteristicas Extraidas (Features)",
        ["Categoria", "Variaveis", "Descricao"],
        [
            ["Temporal", "persistence_score, hour, month, dayofyear", "Padroes temporais da deteccao"],
            ["Sensor", "brightness, frp, confidence, scan, track", "Medidas do sensor satelite"],
            ["Meteorologica", "temperature, humidity, wind, precipitation", "Condicoes climaticas locais"],
            ["Vegetacao", "NDVI, NBR, dNBR pre/pos", "Estado da cobertura vegetal"],
            ["Espacial", "land_cover, distance_urban, elevation", "Contexto geografico"]
        ])

    add_content_slide(prs, "Processo de Weak Labeling", [
        "Para cada deteccao FIRMS:",
        "",
        "1. Localizar pixels MCD64A1 em buffer de 1km",
        "2. Verificar data de queima na janela de +/- 15 dias",
        "3. Atribuir rotulo binario:",
        "",
        "  REAL (label=1):",
        "    Existe cicatriz de queimada proxima no periodo",
        "",
        "  ESPURIA (label=0):",
        "    Nenhuma correspondencia encontrada",
        "",
        "Dataset resultante: 708 deteccoes rotuladas"
    ], LARANJA)

    add_two_column_slide(prs, "Estrategias de Validacao",
        "Validacao Espacial",
        [
            "Leave-One-State-Out Cross-Validation",
            "",
            "Treina em 3 estados",
            "Testa no 4o estado",
            "",
            "Objetivo:",
            "Avaliar generalizacao geografica",
            "do modelo para areas nao vistas"
        ],
        "Validacao Temporal",
        [
            "Holdout Temporal",
            "",
            "Treino: 2022-2023",
            "Teste: 2024",
            "",
            "Objetivo:",
            "Avaliar estabilidade do modelo",
            "ao longo do tempo (concept drift)"
        ])

    # === SECAO 5: RESULTADOS ===
    add_section_slide(prs, 5, "Resultados\nExperimentais")

    add_metrics_slide(prs, "Desempenho do Classificador LightGBM",
        [("Acuracia", "82%"), ("ROC-AUC", "86%"), ("PR-AUC", "81%"), ("F1-Score", "80%")])

    add_table_slide(prs, "Comparacao: LightGBM vs XGBoost",
        ["Metrica", "LightGBM", "XGBoost", "Diferenca"],
        [
            ["Acuracia", "82%", "81%", "+1 p.p."],
            ["ROC-AUC", "86%", "85%", "+1 p.p."],
            ["PR-AUC", "81%", "80%", "+1 p.p."],
            ["Tempo de Treino", "0.8s", "2.1s", "2.6x mais rapido"]
        ])

    add_table_slide(prs, "Validacao Espacial: Leave-One-State-Out",
        ["Estado (Teste)", "Acuracia", "ROC-AUC", "PR-AUC"],
        [
            ["Maranhao", "80%", "84%", "79%"],
            ["Tocantins", "83%", "87%", "82%"],
            ["Piaui", "81%", "85%", "80%"],
            ["Bahia", "82%", "86%", "81%"]
        ])

    add_content_slide(prs, "Validacao Temporal: Resultados", [
        "Treino (2022-2023) -> Teste (2024):",
        "",
        "  Acuracia: 82% -> 77%  (queda de 5 p.p.)",
        "  ROC-AUC: 86% -> 81%  (queda de 5 p.p.)",
        "  PR-AUC: 81% -> 76%  (queda de 5 p.p.)",
        "",
        "Interpretacao:",
        "  Indica concept drift moderado ao longo do tempo",
        "  Padroes de queimadas podem variar entre anos",
        "  Recomendacao: Retreinamento periodico do modelo"
    ], CYAN)

    add_content_slide(prs, "Importancia das Caracteristicas", [
        "Top 5 Features por Importancia (Gain):",
        "",
        "1. persistence_score  (35%)",
        "  Focos reais tendem a persistir por mais tempo",
        "2. ndvi  (18%)",
        "  Vegetacao presente indica combustivel disponivel",
        "3. frp - Fire Radiative Power  (12%)",
        "  Potencia radiativa maior em incendios reais",
        "4. temperature_2m  (9%)",
        "  Condicoes meteorologicas influenciam propagacao",
        "5. brightness  (7%)",
        "  Temperatura de brilho do pixel"
    ], CYAN)

    add_metrics_slide(prs, "Impacto Operacional Estimado",
        [("Reducao de FP", "87%"), ("Alertas Filtrados", "13/100"), ("Ganho de Eficiencia", "7x"), ("Recall Mantido", "95%")])

    # === SECAO 6: CRONOGRAMA ===
    add_section_slide(prs, 6, "Cronograma")

    add_two_column_slide(prs, "Cronograma de Atividades",
        "Concluido",
        [
            "Revisao bibliografica sistematica",
            "Implementacao do pipeline de dados",
            "Desenvolvimento do weak labeling",
            "Extracao de features",
            "Treinamento dos modelos",
            "Validacao espacial e temporal",
            "Analise de importancia de features"
        ],
        "Pendente",
        [
            "Expansao do dataset (mais anos)",
            "Experimentos com deep learning",
            "Integracao com MapBiomas Fire",
            "Modulo de predicao D+1",
            "Redacao final da dissertacao",
            "Defesa da dissertacao"
        ])

    # === SECAO 7: CONCLUSAO ===
    add_section_slide(prs, 7, "Conclusoes")

    add_content_slide(prs, "Contribuicoes do Trabalho", [
        "1. Pipeline de dados geoespaciais multifonte",
        "  Integracao automatizada FIRMS + MCD64A1 + ERA5 + Sentinel-2",
        "",
        "2. Metodologia de weak labeling para focos de incendio",
        "  Rotulacao automatica escalavel usando MCD64A1",
        "",
        "3. Conjunto de features para classificacao",
        "  20+ caracteristicas em 5 categorias",
        "",
        "4. Modelo validado espacial e temporalmente",
        "  82% acuracia, 86% ROC-AUC com LightGBM",
        "",
        "5. Analise comparativa de algoritmos"
    ], VERDE)

    add_content_slide(prs, "Limitacoes e Trabalhos Futuros", [
        "Limitacoes identificadas:",
        "  Dependencia do MCD64A1 (resolucao e latencia)",
        "  Dataset relativamente pequeno (708 amostras)",
        "  Especificidade regional (apenas MATOPIBA)",
        "",
        "Trabalhos futuros:",
        "  Expansao do dataset temporal e espacial",
        "  Integracao com MapBiomas Fire (maior resolucao)",
        "  Experimentos com arquiteturas de deep learning",
        "  Modulo de predicao de propagacao (D+1)",
        "  Sistema operacional de monitoramento"
    ], LARANJA)

    # SLIDE FINAL
    add_end_slide(prs)

    prs.save("Apresentacao_Qualificacao_Mestrado.pptx")
    print("Apresentacao salva: Apresentacao_Qualificacao_Mestrado.pptx")


def create_roteiro():
    """Cria roteiro em DOCX"""
    doc = Document()

    title = doc.add_heading("Roteiro de Apresentacao - Qualificacao de Mestrado", 0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER

    doc.add_paragraph("Classificacao Automatica de Deteccoes de Focos de Incendio na Regiao MATOPIBA")
    doc.add_paragraph("Tempo total: 20-25 minutos")
    doc.add_paragraph("=" * 60)

    roteiro = [
        ("SLIDE 1: Titulo", "30s", [
            "Cumprimentar a banca e agradecer a presenca",
            "Apresentar-se: nome, programa, orientadores",
            "Mencionar o titulo do trabalho"
        ]),
        ("SLIDE 2: Agenda", "30s", [
            "Apresentar a estrutura da apresentacao",
            "Mencionar que perguntas serao ao final"
        ]),
        ("SLIDE 3: Secao - Introducao", "10s", [
            "Transicao: 'Vamos comecar pelo contexto do problema...'"
        ]),
        ("SLIDE 4: Contexto", "1.5min", [
            "Brasil tem alta incidencia de incendios florestais",
            "INPE monitora via satelite - sistema FIRMS",
            "Problema CENTRAL: muitos falsos positivos",
            "Isso sobrecarrega os orgaos de monitoramento",
            "ENFATIZAR: verificacao manual e inviavel"
        ]),
        ("SLIDE 5: O Problema", "1min", [
            "Listar as fontes de falsos positivos",
            "Explicar o impacto operacional",
            "Motivar a necessidade de automacao"
        ]),
        ("SLIDE 6: MATOPIBA", "1min", [
            "Mostrar a regiao geograficamente",
            "Explicar por que foi escolhida",
            "Destacar: disponibilidade de dados"
        ]),
        ("SLIDE 7: Secao - Objetivos", "10s", [
            "Transicao para objetivos"
        ]),
        ("SLIDE 8: Objetivo Geral", "30s", [
            "Ler de forma clara e pausada",
            "ENFATIZAR: distinguir REAL de ESPURIA",
            "Mencionar: machine learning com LightGBM"
        ]),
        ("SLIDE 9: Objetivos Especificos", "1min", [
            "Passar por cada objetivo rapidamente",
            "Destacar que TODOS foram alcancados"
        ]),
        ("SLIDE 10: Secao - Fundamentacao", "10s", [
            "Transicao para base teorica"
        ]),
        ("SLIDE 11: Sensoriamento Remoto", "1.5min", [
            "Explicar deteccao por anomalia termica",
            "Apresentar MODIS e VIIRS",
            "Descrever os produtos de dados usados"
        ]),
        ("SLIDE 12: Gradient Boosting", "1min", [
            "Explicar de forma SIMPLES - ensemble de arvores",
            "Cada arvore corrige erros anteriores",
            "Mencionar LightGBM como escolha principal"
        ]),
        ("SLIDE 13: Weak Labeling", "1min", [
            "IMPORTANTE: explicar bem este conceito",
            "Problema: rotular manualmente e impossivel",
            "Solucao: usar MCD64A1 como referencia",
            "Detalhar janela de +/- 15 dias"
        ]),
        ("SLIDE 14: Secao - Metodologia", "10s", [
            "Transicao para metodologia"
        ]),
        ("SLIDE 15: Arquitetura", "1min", [
            "Apresentar as 4 etapas do pipeline",
            "Destacar modularidade do sistema"
        ]),
        ("SLIDE 16: Fontes de Dados", "1min", [
            "Apresentar cada fonte brevemente",
            "Destacar integracao multifonte"
        ]),
        ("SLIDE 17: Features", "1min", [
            "Apresentar categorias de features",
            "Mencionar: 20+ caracteristicas extraidas"
        ]),
        ("SLIDE 18: Weak Labeling", "1min", [
            "Detalhar o processo de rotulacao",
            "Resultado: 708 amostras rotuladas"
        ]),
        ("SLIDE 19: Validacao", "1min", [
            "Explicar leave-one-state-out",
            "Explicar holdout temporal",
            "DESTACAR: rigor metodologico"
        ]),
        ("SLIDE 20: Secao - Resultados", "10s", [
            "Transicao: 'Agora os resultados principais...'"
        ]),
        ("SLIDE 21: Metricas", "1min", [
            "MEMORIZAR: 82% acuracia, 86% ROC-AUC",
            "Destacar que sao resultados solidos",
            "Comparar com trabalhos relacionados se possivel"
        ]),
        ("SLIDE 22: LightGBM vs XGBoost", "30s", [
            "LightGBM ligeiramente superior",
            "DESTAQUE: 2.6x mais rapido para treinar"
        ]),
        ("SLIDE 23: Validacao Espacial", "1min", [
            "Mostrar consistencia entre estados",
            "Variacao < 4 pontos percentuais",
            "BOA generalizacao geografica"
        ]),
        ("SLIDE 24: Validacao Temporal", "1min", [
            "Mostrar queda de ~5 p.p.",
            "Explicar concept drift",
            "Mencionar necessidade de retreinamento"
        ]),
        ("SLIDE 25: Features", "1min", [
            "persistence_score e a mais importante",
            "Explicar POR QUE faz sentido",
            "NDVI indica combustivel disponivel"
        ]),
        ("SLIDE 26: Impacto", "1min", [
            "DESTACAR: reducao de 87% em falsos positivos",
            "De 100 alertas, apenas 13 precisam verificacao",
            "Recall de 95% - quase nenhum incendio real perdido"
        ]),
        ("SLIDE 27: Secao - Cronograma", "10s", [
            "Transicao para cronograma"
        ]),
        ("SLIDE 28: Cronograma", "30s", [
            "Mostrar o que foi concluido",
            "Apresentar proximas etapas"
        ]),
        ("SLIDE 29: Secao - Conclusoes", "10s", [
            "Transicao final"
        ]),
        ("SLIDE 30: Contribuicoes", "1min", [
            "Enumerar as 5 contribuicoes",
            "DESTACAR originalidade do trabalho"
        ]),
        ("SLIDE 31: Limitacoes", "30s", [
            "Ser HONESTO sobre limitacoes",
            "Mostra maturidade academica"
        ]),
        ("SLIDE 32: Obrigado", "30s", [
            "Agradecer novamente a banca",
            "Abrir para perguntas"
        ])
    ]

    for titulo, tempo, pontos in roteiro:
        doc.add_heading(titulo, level=1)
        doc.add_paragraph(f"Tempo estimado: {tempo}")
        for p in pontos:
            doc.add_paragraph(p, style='List Bullet')
        doc.add_paragraph("")

    doc.add_heading("PERGUNTAS PROVAVEIS DA BANCA", level=1)
    perguntas = [
        "Por que escolheu LightGBM?",
        "Qual a confiabilidade do weak labeling?",
        "Como lidou com desbalanceamento de classes?",
        "O sistema funcionaria em outras regioes?",
        "Por que nao usou deep learning?",
        "Como pretende lidar com concept drift em producao?",
        "Quais as limitacoes da resolucao do MCD64A1?",
        "Como validar que os rotulos estao corretos?"
    ]
    for p in perguntas:
        doc.add_paragraph(p, style='List Bullet')

    doc.save("Roteiro_Apresentacao_Qualificacao.docx")
    print("Roteiro salvo: Roteiro_Apresentacao_Qualificacao.docx")


if __name__ == "__main__":
    print("Criando apresentacao moderna...")
    create_presentation()
    print("\nCriando roteiro...")
    create_roteiro()
    print("\nConcluido!")
