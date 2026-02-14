# -*- coding: utf-8 -*-
"""Apresentacao v4 - Com imagens e explicacoes visuais"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# Paleta
AZUL_DARK = RGBColor(0x0d, 0x1b, 0x2a)
AZUL_MEDIO = RGBColor(0x1b, 0x26, 0x3b)
AZUL_ACCENT = RGBColor(0x41, 0x5a, 0x77)
AZUL_LIGHT = RGBColor(0x77, 0x8d, 0xa9)
CYAN = RGBColor(0x00, 0xd4, 0xff)
LARANJA = RGBColor(0xff, 0x6b, 0x35)
VERDE = RGBColor(0x00, 0xc9, 0xa7)
AMARELO = RGBColor(0xff, 0xc8, 0x00)
BRANCO = RGBColor(0xff, 0xff, 0xff)
CINZA_CLARO = RGBColor(0xe0, 0xe1, 0xdd)

# Paths
LOGO_UFAM = "logos/ufam.png"
LOGO_PPGEE = "logos/ppge.png"
LOGO_POLITECNICO = "logos/politecnico de portalegre.png"

IMG_DISTRIBUICAO = "docs/visual/01_distribuicao_hotspots.png"
IMG_PREDICOES = "docs/visual/02_predicoes_modelo.png"
IMG_ESTRATEGIAS = "docs/visual/03_estrategias_filtro.png"
IMG_RECOMENDACOES = "docs/visual/04_mapa_recomendacoes.png"
IMG_WORKFLOW = "docs/visual/05_workflow_diagram.png"

def set_gradient(shape, color1, color2, angle=90):
    fill = shape.fill
    fill.gradient()
    fill.gradient_angle = angle
    fill.gradient_stops[0].color.rgb = color1
    fill.gradient_stops[1].color.rgb = color2

def add_decorations(slide, prs, style="content"):
    if style == "title":
        tri = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Inches(9.5), Inches(-0.5), Inches(4.5), Inches(5))
        tri.fill.solid()
        tri.fill.fore_color.rgb = RGBColor(0x15, 0x20, 0x30)
        tri.line.fill.background()
        tri.rotation = 180

        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10), Inches(4.5), Inches(4), Inches(4))
        circ.fill.solid()
        circ.fill.fore_color.rgb = CYAN
        circ.fill.fore_color.brightness = 0.7
        circ.line.fill.background()

    elif style == "section":
        hex1 = slide.shapes.add_shape(MSO_SHAPE.HEXAGON, Inches(10), Inches(-1), Inches(5), Inches(5))
        hex1.fill.solid()
        hex1.fill.fore_color.rgb = CYAN
        hex1.fill.fore_color.brightness = 0.8
        hex1.line.fill.background()

    elif style == "content":
        circ1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(12.5), Inches(-0.3), Inches(1.2), Inches(1.2))
        circ1.fill.solid()
        circ1.fill.fore_color.rgb = CYAN
        circ1.fill.fore_color.brightness = 0.85
        circ1.line.fill.background()

        circ2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(12.7), Inches(0.7), Inches(0.7), Inches(0.7))
        circ2.fill.solid()
        circ2.fill.fore_color.rgb = VERDE
        circ2.fill.fore_color.brightness = 0.85
        circ2.line.fill.background()

def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    set_gradient(bg, AZUL_DARK, AZUL_MEDIO, 135)
    bg.line.fill.background()

    add_decorations(slide, prs, "title")

    # Logo bar
    logo_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    logo_bar.fill.solid()
    logo_bar.fill.fore_color.rgb = BRANCO
    logo_bar.line.fill.background()

    if os.path.exists(LOGO_UFAM):
        slide.shapes.add_picture(LOGO_UFAM, Inches(0.3), Inches(0.1), height=Inches(1))
    if os.path.exists(LOGO_PPGEE):
        slide.shapes.add_picture(LOGO_PPGEE, Inches(5.5), Inches(0.15), height=Inches(0.9))
    if os.path.exists(LOGO_POLITECNICO):
        slide.shapes.add_picture(LOGO_POLITECNICO, Inches(10), Inches(0.25), height=Inches(0.7))

    colors = [CYAN, VERDE, LARANJA, AMARELO]
    for i, color in enumerate(colors):
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(i * 3.333), Inches(1.2), Inches(3.333), Pt(5))
        line.fill.solid()
        line.fill.fore_color.rgb = color
        line.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(2.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Classificacao Automatica de Deteccoes\nde Focos de Incendio na Regiao\nMATOPIBA Usando Machine Learning"
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = BRANCO
    p.font.name = "Segoe UI"

    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.4), Inches(9), Inches(0.6))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Exame de Qualificacao | Mestrado em Engenharia Eletrica"
    p.font.size = Pt(18)
    p.font.color.rgb = CYAN
    p.font.name = "Segoe UI Light"

    info_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.3), Inches(8), Inches(1.8))
    tf = info_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Mestrando: Andrey Ruben Ribeiro Bessa"
    p.font.size = Pt(14)
    p.font.color.rgb = CINZA_CLARO
    p.font.name = "Segoe UI"
    p = tf.add_paragraph()
    p.text = "Orientador: Prof. D.Sc. Celso Barbosa Carvalho"
    p.font.size = Pt(14)
    p.font.color.rgb = CINZA_CLARO
    p.font.name = "Segoe UI"
    p = tf.add_paragraph()
    p.text = "Coorientador: Prof. D.Sc. Andre Chaves Mendes"
    p.font.size = Pt(14)
    p.font.color.rgb = CINZA_CLARO
    p.font.name = "Segoe UI"

def add_section_slide(prs, number, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    set_gradient(bg, AZUL_DARK, AZUL_MEDIO, 45)
    bg.line.fill.background()

    add_decorations(slide, prs, "section")

    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1), Inches(2.2), Inches(2.8), Inches(2.8))
    circle.fill.solid()
    circle.fill.fore_color.rgb = CYAN
    circle.line.fill.background()

    ring = slide.shapes.add_shape(MSO_SHAPE.DONUT, Inches(0.7), Inches(1.9), Inches(3.4), Inches(3.4))
    ring.fill.solid()
    ring.fill.fore_color.rgb = VERDE
    ring.fill.fore_color.brightness = 0.3
    ring.line.fill.background()

    num_box = slide.shapes.add_textbox(Inches(1), Inches(2.7), Inches(2.8), Inches(2))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{number:02d}"
    p.font.size = Pt(80)
    p.font.bold = True
    p.font.color.rgb = AZUL_DARK
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Segoe UI Black"

    title_box = slide.shapes.add_textbox(Inches(4.5), Inches(2.5), Inches(8), Inches(2.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = BRANCO
    p.font.name = "Segoe UI"

def add_content_slide(prs, title, bullets, accent_color=CYAN):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BRANCO
    bg.line.fill.background()

    sidebar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.2), prs.slide_height)
    set_gradient(sidebar, accent_color, AZUL_DARK, 180)
    sidebar.line.fill.background()

    add_decorations(slide, prs, "content")

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(11), Inches(0.9))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = AZUL_DARK
    p.font.name = "Segoe UI"

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.15), Inches(1.5), Pt(4))
    line.fill.solid()
    line.fill.fore_color.rgb = accent_color
    line.line.fill.background()

    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(11.5), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if bullet.startswith("  "):
            p.text = "     " + bullet.strip()
            p.font.size = Pt(17)
            p.font.color.rgb = AZUL_ACCENT
        elif bullet == "":
            p.text = ""
            p.font.size = Pt(10)
        else:
            p.text = bullet
            p.font.size = Pt(19)
            p.font.color.rgb = AZUL_DARK
        p.font.name = "Segoe UI"
        p.space_after = Pt(8)

def add_image_slide(prs, title, image_path, explanations, accent_color=CYAN):
    """Slide com imagem e explicacoes ao lado"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BRANCO
    bg.line.fill.background()

    sidebar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.2), prs.slide_height)
    set_gradient(sidebar, accent_color, AZUL_DARK, 180)
    sidebar.line.fill.background()

    # Titulo
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = AZUL_DARK
    p.font.name = "Segoe UI"

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.9), Inches(1.2), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = accent_color
    line.line.fill.background()

    # Imagem
    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(0.4), Inches(1.15), width=Inches(7.8))

    # Explicacoes ao lado
    exp_box = slide.shapes.add_textbox(Inches(8.5), Inches(1.15), Inches(4.5), Inches(6))
    tf = exp_box.text_frame
    tf.word_wrap = True

    for i, exp in enumerate(explanations):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if exp.startswith(">>"):
            p.text = exp[2:].strip()
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = accent_color
        else:
            p.text = exp
            p.font.size = Pt(13)
            p.font.color.rgb = AZUL_DARK
        p.font.name = "Segoe UI"
        p.space_after = Pt(6)

    return slide

def add_image_fullwidth_slide(prs, title, image_path, caption, accent_color=CYAN):
    """Slide com imagem em largura total e legenda embaixo"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BRANCO
    bg.line.fill.background()

    sidebar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.2), prs.slide_height)
    set_gradient(sidebar, accent_color, AZUL_DARK, 180)
    sidebar.line.fill.background()

    # Titulo
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = AZUL_DARK
    p.font.name = "Segoe UI"

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.9), Inches(1.2), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = accent_color
    line.line.fill.background()

    # Imagem centralizada
    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(0.5), Inches(1.1), width=Inches(12.3))

    # Legenda
    cap_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.8))
    tf = cap_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = caption
    p.font.size = Pt(14)
    p.font.color.rgb = AZUL_ACCENT
    p.font.name = "Segoe UI"
    p.alignment = PP_ALIGN.CENTER

    return slide

def add_metrics_slide(prs, title, metrics):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    set_gradient(bg, AZUL_DARK, AZUL_MEDIO, 180)
    bg.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12), Inches(0.9))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = BRANCO
    p.font.name = "Segoe UI"

    colors = [CYAN, VERDE, LARANJA, AMARELO]
    card_width = Inches(2.9)
    card_height = Inches(3.2)
    start_x = Inches(0.5)
    y_pos = Inches(1.8)
    spacing = Inches(3.15)

    for i, (name, value) in enumerate(metrics):
        x = start_x + (i * spacing)
        color = colors[i % 4]

        shadow = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.08), y_pos + Inches(0.08), card_width, card_height)
        shadow.fill.solid()
        shadow.fill.fore_color.rgb = RGBColor(0x05, 0x0a, 0x15)
        shadow.line.fill.background()

        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y_pos, card_width, card_height)
        card.fill.solid()
        card.fill.fore_color.rgb = AZUL_MEDIO
        card.line.color.rgb = color
        card.line.width = Pt(3)

        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y_pos, card_width, Inches(0.15))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()

        icon = slide.shapes.add_shape(MSO_SHAPE.HEXAGON, x + Inches(0.9), y_pos + Inches(0.4), Inches(1.1), Inches(1.1))
        icon.fill.solid()
        icon.fill.fore_color.rgb = color
        icon.fill.fore_color.brightness = 0.3
        icon.line.fill.background()

        val_box = slide.shapes.add_textbox(x, y_pos + Inches(1.6), card_width, Inches(0.9))
        tf = val_box.text_frame
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER
        p.font.name = "Segoe UI Black"

        name_box = slide.shapes.add_textbox(x, y_pos + Inches(2.4), card_width, Inches(0.8))
        tf = name_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(15)
        p.font.color.rgb = CINZA_CLARO
        p.alignment = PP_ALIGN.CENTER
        p.font.name = "Segoe UI"

def add_table_slide(prs, title, headers, rows):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BRANCO
    bg.line.fill.background()

    sidebar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.2), prs.slide_height)
    set_gradient(sidebar, CYAN, VERDE, 180)
    sidebar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(11), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = AZUL_DARK
    p.font.name = "Segoe UI"

    num_cols = len(headers)
    num_rows = len(rows) + 1
    table = slide.shapes.add_table(num_rows, num_cols, Inches(0.5), Inches(1.4), Inches(12), Inches(0.55 * num_rows)).table

    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = AZUL_DARK
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(13)
        p.font.color.rgb = BRANCO
        p.alignment = PP_ALIGN.CENTER
        p.font.name = "Segoe UI"

    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = str(val)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            p.font.color.rgb = AZUL_DARK
            p.alignment = PP_ALIGN.CENTER
            p.font.name = "Segoe UI"
            if ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xf0, 0xf8, 0xff)

def add_two_column_slide(prs, title, left_title, left_items, right_title, right_items):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BRANCO
    bg.line.fill.background()

    sidebar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.2), prs.slide_height)
    set_gradient(sidebar, VERDE, CYAN, 180)
    sidebar.line.fill.background()

    divider = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.4), Inches(1.5), Inches(0.1), Inches(5.5))
    set_gradient(divider, CYAN, VERDE, 180)
    divider.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = AZUL_DARK
    p.font.name = "Segoe UI"

    # Left
    icon_left = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), Inches(1.3), Inches(0.5), Inches(0.5))
    icon_left.fill.solid()
    icon_left.fill.fore_color.rgb = CYAN
    icon_left.line.fill.background()

    left_title_box = slide.shapes.add_textbox(Inches(1.1), Inches(1.35), Inches(5), Inches(0.5))
    tf = left_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = CYAN
    p.font.name = "Segoe UI"

    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(5.7), Inches(5))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(15)
        p.font.color.rgb = AZUL_DARK
        p.font.name = "Segoe UI"
        p.space_after = Pt(6)

    # Right
    icon_right = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.7), Inches(1.3), Inches(0.5), Inches(0.5))
    icon_right.fill.solid()
    icon_right.fill.fore_color.rgb = LARANJA
    icon_right.line.fill.background()

    right_title_box = slide.shapes.add_textbox(Inches(7.3), Inches(1.35), Inches(5), Inches(0.5))
    tf = right_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = LARANJA
    p.font.name = "Segoe UI"

    right_box = slide.shapes.add_textbox(Inches(6.7), Inches(2), Inches(5.7), Inches(5))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(right_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(15)
        p.font.color.rgb = AZUL_DARK
        p.font.name = "Segoe UI"
        p.space_after = Pt(6)

def add_end_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    set_gradient(bg, AZUL_DARK, AZUL_MEDIO, 135)
    bg.line.fill.background()

    shapes_data = [
        (MSO_SHAPE.HEXAGON, 0.5, 0.5, 2, 2, CYAN, 0.8),
        (MSO_SHAPE.PENTAGON, 11, 0.3, 1.5, 1.5, VERDE, 0.85),
        (MSO_SHAPE.OVAL, 1.5, 5.5, 1.8, 1.8, LARANJA, 0.8),
        (MSO_SHAPE.DIAMOND, 10.5, 5, 2, 2, AMARELO, 0.85),
    ]
    for shape_type, x, y, w, h, color, brightness in shapes_data:
        s = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
        s.fill.solid()
        s.fill.fore_color.rgb = color
        s.fill.fore_color.brightness = brightness
        s.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0), Inches(2.3), prs.slide_width, Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Obrigado!"
    p.font.size = Pt(80)
    p.font.bold = True
    p.font.color.rgb = BRANCO
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Segoe UI"

    sub_box = slide.shapes.add_textbox(Inches(0), Inches(4), prs.slide_width, Inches(0.8))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Perguntas?"
    p.font.size = Pt(32)
    p.font.color.rgb = CYAN
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Segoe UI Light"

    logo_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(6.3), prs.slide_width, Inches(1.2))
    logo_bar.fill.solid()
    logo_bar.fill.fore_color.rgb = BRANCO
    logo_bar.fill.fore_color.brightness = 0.95
    logo_bar.line.fill.background()

    if os.path.exists(LOGO_UFAM):
        slide.shapes.add_picture(LOGO_UFAM, Inches(0.5), Inches(6.4), height=Inches(1))
    if os.path.exists(LOGO_PPGEE):
        slide.shapes.add_picture(LOGO_PPGEE, Inches(5.5), Inches(6.5), height=Inches(0.8))
    if os.path.exists(LOGO_POLITECNICO):
        slide.shapes.add_picture(LOGO_POLITECNICO, Inches(10.5), Inches(6.55), height=Inches(0.6))


def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # SLIDE 1: Titulo
    add_title_slide(prs)

    # SLIDE 2: Agenda
    add_content_slide(prs, "Agenda", [
        "01  Introducao e Motivacao",
        "02  Objetivos do Trabalho",
        "03  Fundamentacao Teorica",
        "04  Metodologia Proposta",
        "05  Resultados Experimentais",
        "06  Cronograma",
        "07  Conclusoes e Proximos Passos"
    ])

    # === SECAO 1: INTRODUCAO ===
    add_section_slide(prs, 1, "Introducao e\nMotivacao")

    add_content_slide(prs, "Contexto: Incendios Florestais no Brasil", [
        "Brasil: um dos paises com maior incidencia de incendios florestais",
        "INPE monitora milhares de focos de calor diariamente via satelite",
        "Sistema FIRMS/NASA gera alertas automaticos em tempo quase-real",
        "  Deteccoes baseadas em anomalias termicas (MODIS, VIIRS)",
        "",
        "Problema critico: Alta taxa de FALSOS POSITIVOS",
        "  Sobrecarrega orgaos de monitoramento ambiental",
        "  Verificacao manual inviavel em escala operacional"
    ], CYAN)

    add_content_slide(prs, "O Problema: Deteccoes Espurias", [
        "Fontes comuns de falsos positivos:",
        "  Reflexos solares em superficies metalicas",
        "  Queima controlada de gas em industrias",
        "  Alvos urbanos com alta temperatura aparente",
        "",
        "Impacto operacional:",
        "  Desperdicio de recursos humanos e financeiros",
        "  Atrasos na resposta a incendios reais",
        "  Perda de credibilidade do sistema de alertas"
    ], LARANJA)

    # === SECAO 2: OBJETIVOS ===
    add_section_slide(prs, 2, "Objetivos")

    add_content_slide(prs, "Objetivo Geral", [
        "Desenvolver um sistema de classificacao automatica",
        "baseado em aprendizado de maquina para:",
        "",
        "  Distinguir deteccoes de incendios REAIS de ESPURIAS",
        "  nos dados do sistema FIRMS/NASA",
        "",
        "Aplicacao: Regiao MATOPIBA (2022-2024)",
        "Algoritmo principal: LightGBM (Gradient Boosting)"
    ], CYAN)

    add_content_slide(prs, "Objetivos Especificos", [
        "1. Implementar pipeline de ingestao de dados geoespaciais",
        "2. Desenvolver metodologia de weak labeling automatica",
        "3. Definir e extrair caracteristicas para classificacao",
        "4. Treinar e otimizar modelos de gradient boosting",
        "5. Validar com estrategias espaciais e temporais",
        "6. Analisar importancia das caracteristicas"
    ], CYAN)

    # === SECAO 3: FUNDAMENTACAO ===
    add_section_slide(prs, 3, "Fundamentacao\nTeorica")

    add_content_slide(prs, "Sensoriamento Remoto: Deteccao de Incendios", [
        "Principio: Deteccao de anomalias termicas via infravermelho",
        "",
        "Satelites utilizados:",
        "  MODIS (Terra/Aqua) - resolucao 1km, 4x ao dia",
        "  VIIRS (Suomi NPP/NOAA-20) - resolucao 375m",
        "",
        "Produtos de dados:",
        "  FIRMS: Focos ativos em tempo quase-real",
        "  MCD64A1: Area queimada mensal (500m)"
    ], VERDE)

    add_content_slide(prs, "Machine Learning: Gradient Boosting", [
        "Tecnica de ensemble baseada em arvores de decisao",
        "",
        "Funcionamento:",
        "  Arvores treinadas sequencialmente",
        "  Cada arvore corrige erros das anteriores",
        "",
        "Algoritmos utilizados:",
        "  LightGBM: Otimizado para velocidade e memoria",
        "  XGBoost: Referencia em competicoes de ML"
    ], VERDE)

    # === SECAO 4: METODOLOGIA ===
    add_section_slide(prs, 4, "Metodologia")

    # WORKFLOW - Imagem com explicacao
    add_image_fullwidth_slide(prs,
        "Arquitetura do Sistema: Pipeline de Classificacao",
        IMG_WORKFLOW,
        "Fluxo completo: INPUT (CSV) -> MODELO (LightGBM) -> OUTPUT (Predicoes) -> FILTRO (USE/REVISAR/DESCARTE)",
        LARANJA)

    add_content_slide(prs, "Detalhamento do Pipeline", [
        "1. INPUT: CSV com novos hotspots",
        "  latitude, longitude, confidence, acq_datetime",
        "",
        "2. MODELO: LightGBM (ou XGBoost)",
        "  Extrai features -> Faz predicoes -> Calcula probabilidade",
        "",
        "3. OUTPUT: Predicoes com confianca",
        "  prediction, spurious_probability, confidence_score",
        "",
        "4. OPCOES DE FILTRO:",
        "  USE (prob < 0.5): 158 hotspots - usar diretamente",
        "  REVISAR (0.5 <= prob < 0.7): 21 hotspots - verificar",
        "  DESCARTE (prob >= 0.7): 51 hotspots - descartar"
    ], LARANJA)

    add_table_slide(prs, "Fontes de Dados Integradas",
        ["Fonte", "Tipo", "Resolucao", "Frequencia"],
        [
            ["FIRMS/VIIRS", "Focos de calor", "375m", "~3 horas"],
            ["MCD64A1", "Area queimada", "500m", "Mensal"],
            ["Sentinel-2", "Multiespectral", "10-20m", "5 dias"],
            ["ERA5-Land", "Meteorologia", "~9km", "Horario"]
        ])

    add_two_column_slide(prs, "Estrategias de Validacao",
        "Validacao Espacial",
        ["Leave-One-State-Out", "", "Treina em 3 estados", "Testa no 4o estado", "", "Avalia generalizacao", "geografica"],
        "Validacao Temporal",
        ["Holdout Temporal", "", "Treino: 2022-2023", "Teste: 2024", "", "Avalia estabilidade", "ao longo do tempo"])

    # === SECAO 5: RESULTADOS ===
    add_section_slide(prs, 5, "Resultados\nExperimentais")

    # IMAGEM 1: Distribuicao dos Hotspots
    add_image_slide(prs,
        "Distribuicao dos Hotspots: Real vs Espurio",
        IMG_DISTRIBUICAO,
        [
            ">> Mapa de Distribuicao",
            "Hotspots Reais (vermelho)",
            "concentrados em clusters",
            "",
            "Hotspots Espurios (verde)",
            "dispersos na regiao",
            "",
            ">> Distribuicao de Confianca",
            "Hotspots Reais: alta",
            "confianca FIRMS (~85-95)",
            "",
            "Hotspots Espurios: maior",
            "variabilidade na confianca",
            "",
            ">> Insight Principal",
            "Confianca FIRMS sozinha",
            "nao e suficiente para",
            "distinguir os tipos"
        ], CYAN)

    # IMAGEM 2: Predicoes do Modelo
    add_image_slide(prs,
        "Desempenho do Modelo LightGBM",
        IMG_PREDICOES,
        [
            ">> Predicoes no Mapa",
            "Predito Real (verde)",
            "Predito Espurio (vermelho X)",
            "",
            ">> Matriz de Confusao",
            "8 (3.5%) - Falsos Negativos",
            "72 (31.3%) - Verdadeiros Neg.",
            "0 (0.0%) - Falsos Positivos",
            "",
            ">> Curva de Probabilidade",
            "Threshold em 0.5 (linha azul)",
            "Alta confianca em 0.7 (laranja)",
            "",
            ">> Distribuicao Probabilidades",
            "Separacao clara entre",
            "classes Real e Espuria"
        ], VERDE)

    add_metrics_slide(prs, "Metricas de Desempenho: LightGBM",
        [("Acuracia", "82%"), ("ROC-AUC", "86%"), ("PR-AUC", "81%"), ("F1-Score", "80%")])

    add_table_slide(prs, "Comparacao: LightGBM vs XGBoost",
        ["Metrica", "LightGBM", "XGBoost", "Vantagem"],
        [
            ["Acuracia", "82%", "81%", "+1 p.p."],
            ["ROC-AUC", "86%", "85%", "+1 p.p."],
            ["Tempo Treino", "0.8s", "2.1s", "2.6x mais rapido"]
        ])

    # IMAGEM 3: Estrategias de Filtro
    add_image_slide(prs,
        "Comparacao de Estrategias de Filtro",
        IMG_ESTRATEGIAS,
        [
            ">> Cenario 1: Sem Filtro",
            "Usar TODOS os hotspots",
            "Acertos: 150/230",
            "Taxa: 65.2% (vermelho)",
            "",
            ">> Cenario 2: Predito=Real",
            "Usar apenas preditos Real",
            "Hotspots: 72 (31.3%)",
            "Taxa: 0.0% - NAO FUNCIONA",
            "",
            ">> Cenario 3: Prob < 0.5",
            "Usar probabilidade espuria",
            "Hotspots: 158 (68.7%)",
            "Acertos: 150/158",
            "Taxa: 94.9% (verde)",
            "",
            ">> MELHOR ESTRATEGIA",
            "Filtrar por prob < 0.5"
        ], LARANJA)

    # IMAGEM 4: Mapa de Recomendacoes
    add_image_slide(prs,
        "Mapa de Recomendacoes Operacionais",
        IMG_RECOMENDACOES,
        [
            ">> Legenda do Mapa",
            "",
            "RECOMENDADO (laranja)",
            "n=158 hotspots",
            "Usar na analise",
            "",
            "DESCARTADO (verde)",
            "n=51 hotspots",
            "Alta prob. espuria",
            "",
            "REVISAR (amarelo)",
            "n=21 hotspots",
            "Zona de incerteza",
            "",
            ">> Aplicacao Pratica",
            "Reducao de 87% nos",
            "falsos positivos",
            "mantendo 95% recall"
        ], VERDE)

    add_metrics_slide(prs, "Impacto Operacional Estimado",
        [("Reducao FP", "87%"), ("Alertas", "13/100"), ("Eficiencia", "7x"), ("Recall", "95%")])

    add_table_slide(prs, "Validacao Espacial: Leave-One-State-Out",
        ["Estado (Teste)", "Acuracia", "ROC-AUC", "PR-AUC"],
        [
            ["Maranhao", "80%", "84%", "79%"],
            ["Tocantins", "83%", "87%", "82%"],
            ["Piaui", "81%", "85%", "80%"],
            ["Bahia", "82%", "86%", "81%"]
        ])

    add_content_slide(prs, "Importancia das Caracteristicas", [
        "Top 5 Features por Importancia (Gain):",
        "",
        "1. persistence_score (35%)",
        "  Focos reais persistem por mais tempo",
        "2. ndvi (18%)",
        "  Vegetacao indica combustivel disponivel",
        "3. frp - Fire Radiative Power (12%)",
        "  Potencia maior em incendios reais",
        "4. temperature_2m (9%)",
        "5. brightness (7%)"
    ], CYAN)

    # === SECAO 6: CRONOGRAMA ===
    add_section_slide(prs, 6, "Cronograma")

    add_two_column_slide(prs, "Cronograma de Atividades",
        "Concluido",
        ["Revisao bibliografica", "Pipeline de dados", "Weak labeling", "Extracao de features",
         "Treinamento dos modelos", "Validacao espacial e temporal"],
        "Pendente",
        ["Expansao do dataset", "Deep learning", "MapBiomas Fire", "Predicao D+1",
         "Redacao dissertacao", "Defesa"])

    # === SECAO 7: CONCLUSAO ===
    add_section_slide(prs, 7, "Conclusoes")

    add_content_slide(prs, "Contribuicoes do Trabalho", [
        "1. Pipeline de dados geoespaciais multifonte",
        "  FIRMS + MCD64A1 + ERA5 + Sentinel-2",
        "",
        "2. Metodologia de weak labeling escalavel",
        "",
        "3. Conjunto de 20+ features para classificacao",
        "",
        "4. Modelo validado: 82% acuracia, 86% ROC-AUC",
        "",
        "5. Sistema de filtragem: reducao de 87% em FP"
    ], VERDE)

    add_content_slide(prs, "Limitacoes e Trabalhos Futuros", [
        "Limitacoes:",
        "  Dependencia do MCD64A1 (resolucao/latencia)",
        "  Dataset pequeno (708 amostras)",
        "  Especificidade regional (MATOPIBA)",
        "",
        "Trabalhos futuros:",
        "  Expansao temporal e espacial",
        "  Arquiteturas deep learning",
        "  Modulo predicao propagacao D+1",
        "  Sistema operacional de monitoramento"
    ], LARANJA)

    # SLIDE FINAL
    add_end_slide(prs)

    prs.save("Apresentacao_Qualificacao_v4.pptx")
    print("Apresentacao v4 salva: Apresentacao_Qualificacao_v4.pptx")


if __name__ == "__main__":
    print("Criando apresentacao v4 com imagens...")
    create_presentation()
    print("Concluido!")
