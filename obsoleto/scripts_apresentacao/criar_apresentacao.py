# -*- coding: utf-8 -*-
"""
Script para criar apresentacao PPTX e roteiro DOCX para qualificacao de mestrado
Tema: Classificacao de Focos de Incendio usando Machine Learning na regiao MATOPIBA
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from docx import Document
from docx.shared import Pt as DocxPt, Inches as DocxInches
from docx.enum.text import WD_ALIGN_PARAGRAPH

# Cores do tema
AZUL_ESCURO = RGBColor(0x1a, 0x47, 0x7a)
AZUL_MEDIO = RGBColor(0x2c, 0x6f, 0xb3)
AZUL_CLARO = RGBColor(0x5a, 0x9b, 0xd5)
LARANJA = RGBColor(0xe8, 0x6c, 0x00)
VERDE = RGBColor(0x2e, 0x7d, 0x32)
CINZA_ESCURO = RGBColor(0x33, 0x33, 0x33)
BRANCO = RGBColor(0xff, 0xff, 0xff)

def add_title_slide(prs, title, subtitle):
    """Adiciona slide de titulo"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Fundo azul escuro
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = AZUL_ESCURO
    background.line.fill.background()

    # Linha decorativa
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(4.2), Inches(3), Pt(4))
    line.fill.solid()
    line.fill.fore_color.rgb = LARANJA
    line.line.fill.background()

    # Titulo
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.5), Inches(2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = BRANCO

    # Subtitulo
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12.5), Inches(1.5))
    tf = sub_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(18)
    p.font.color.rgb = BRANCO

    return slide

def add_section_slide(prs, section_title, section_number):
    """Adiciona slide de secao"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Fundo
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = AZUL_MEDIO
    background.line.fill.background()

    # Numero da secao
    num_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(2), Inches(1))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = str(section_number).zfill(2)
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = LARANJA

    # Titulo da secao
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(12.5), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = section_title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = BRANCO

    return slide

def add_content_slide(prs, title, bullet_points):
    """Adiciona slide de conteudo com bullets"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Fundo branco
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = BRANCO
    background.line.fill.background()

    # Barra superior azul
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.1))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = AZUL_ESCURO
    top_bar.line.fill.background()

    # Titulo
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.5), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = AZUL_ESCURO

    # Bullets
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12.5), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, point in enumerate(bullet_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = point
        p.font.size = Pt(20)
        p.font.color.rgb = CINZA_ESCURO
        p.space_after = Pt(12)
        p.level = 0

    return slide

def add_metrics_slide(prs, title, metrics):
    """Adiciona slide com metricas em cards"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Fundo
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = BRANCO
    background.line.fill.background()

    # Barra superior
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.1))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = AZUL_ESCURO
    top_bar.line.fill.background()

    # Titulo
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.5), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = AZUL_ESCURO

    # Cards de metricas
    card_width = Inches(2.8)
    card_height = Inches(2)
    start_x = Inches(0.7)
    y_pos = Inches(2.2)
    spacing = Inches(3.1)

    colors = [AZUL_ESCURO, AZUL_MEDIO, VERDE, LARANJA]

    for i, (metric_name, metric_value) in enumerate(metrics):
        x_pos = start_x + (i * spacing)

        # Card
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, y_pos, card_width, card_height)
        card.fill.solid()
        card.fill.fore_color.rgb = colors[i % len(colors)]
        card.line.fill.background()

        # Valor
        val_box = slide.shapes.add_textbox(x_pos, y_pos + Inches(0.3), card_width, Inches(1))
        tf = val_box.text_frame
        p = tf.paragraphs[0]
        p.text = metric_value
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = BRANCO
        p.alignment = PP_ALIGN.CENTER

        # Nome da metrica
        name_box = slide.shapes.add_textbox(x_pos, y_pos + Inches(1.3), card_width, Inches(0.6))
        tf = name_box.text_frame
        p = tf.paragraphs[0]
        p.text = metric_name
        p.font.size = Pt(16)
        p.font.color.rgb = BRANCO
        p.alignment = PP_ALIGN.CENTER

    return slide

def add_table_slide(prs, title, headers, rows):
    """Adiciona slide com tabela"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Fundo
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = BRANCO
    background.line.fill.background()

    # Barra superior
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.1))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = AZUL_ESCURO
    top_bar.line.fill.background()

    # Titulo
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.5), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = AZUL_ESCURO

    # Tabela
    num_cols = len(headers)
    num_rows = len(rows) + 1
    table_width = Inches(12)
    table_height = Inches(0.5 * num_rows)

    table = slide.shapes.add_table(num_rows, num_cols, Inches(0.7), Inches(1.5), table_width, table_height).table

    # Cabecalho
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = AZUL_ESCURO
        para = cell.text_frame.paragraphs[0]
        para.font.bold = True
        para.font.size = Pt(14)
        para.font.color.rgb = BRANCO
        para.alignment = PP_ALIGN.CENTER

    # Dados
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_text)
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(12)
            para.font.color.rgb = CINZA_ESCURO
            para.alignment = PP_ALIGN.CENTER
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xf5, 0xf5, 0xf5)

    return slide

def create_presentation():
    """Cria a apresentacao completa"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # SLIDE 1: Titulo
    add_title_slide(prs,
        "Classificacao Automatica de Deteccoes de\nFocos de Incendio na Regiao MATOPIBA\nUsando Aprendizado de Maquina",
        "Qualificacao de Mestrado\nPrograma de Pos-Graduacao em Engenharia Eletrica\n\nOrientando: [Seu Nome]\nOrientador: [Nome do Orientador]")

    # SLIDE 2: Agenda
    add_content_slide(prs, "Agenda",
        ["1. Introducao e Motivacao",
         "2. Objetivos",
         "3. Fundamentacao Teorica",
         "4. Metodologia",
         "5. Resultados Experimentais",
         "6. Cronograma",
         "7. Conclusoes e Trabalhos Futuros"])

    # SECAO 1: INTRODUCAO
    add_section_slide(prs, "Introducao e Motivacao", 1)

    # SLIDE 4: Contexto
    add_content_slide(prs, "Contexto: Incendios Florestais no Brasil",
        ["Brasil possui uma das maiores incidencias de incendios florestais do mundo",
         "INPE registra milhares de focos de calor diariamente via satelite",
         "Sistemas de deteccao como FIRMS/NASA geram alertas automaticos",
         "Problema: Alta taxa de falsos positivos sobrecarrega orgaos de monitoramento",
         "Verificacao manual de cada alerta e inviavel em escala"])

    # SLIDE 5: Problema
    add_content_slide(prs, "Problema: Falsos Positivos em Deteccoes",
        ["Fontes de falsos positivos:",
         "  - Reflexos solares em superficies metalicas",
         "  - Queima de gas em industrias",
         "  - Alvos urbanos com alta temperatura",
         "  - Erros de classificacao do sensor",
         "Impacto: Desperdicio de recursos e atrasos na resposta a incendios reais",
         "Necessidade: Sistema automatico para classificar deteccoes"])

    # SLIDE 6: Regiao MATOPIBA
    add_content_slide(prs, "Area de Estudo: Regiao MATOPIBA",
        ["Fronteira agricola no Cerrado brasileiro",
         "Abrange partes de: Maranhao, Tocantins, Piaui e Bahia",
         "Alta incidencia de queimadas (naturais e agricolas)",
         "Relevancia economica e ambiental significativa",
         "Disponibilidade de dados satelitais para o periodo 2022-2024"])

    # SECAO 2: OBJETIVOS
    add_section_slide(prs, "Objetivos", 2)

    # SLIDE 8: Objetivo Geral
    add_content_slide(prs, "Objetivo Geral",
        ["Desenvolver um sistema de classificacao automatica",
         "para distinguir deteccoes de incendios REAIS de ESPURIAS",
         "nos dados do sistema FIRMS/NASA",
         "utilizando tecnicas de aprendizado de maquina"])

    # SLIDE 9: Objetivos Especificos
    add_content_slide(prs, "Objetivos Especificos",
        ["1. Implementar pipeline de ingestao de dados geoespaciais multifonte",
         "2. Desenvolver metodologia de weak labeling automatica",
         "3. Definir e extrair caracteristicas relevantes para classificacao",
         "4. Treinar e otimizar modelos de gradient boosting",
         "5. Validar com estrategias espaciais e temporais rigorosas",
         "6. Analisar importancia das caracteristicas preditivas"])

    # SECAO 3: FUNDAMENTACAO
    add_section_slide(prs, "Fundamentacao Teorica", 3)

    # SLIDE 11: Sensoriamento Remoto
    add_content_slide(prs, "Sensoriamento Remoto para Deteccao de Incendios",
        ["Principio: Deteccao de anomalias termicas via infravermelho",
         "Satelites principais: MODIS (Terra/Aqua), VIIRS (Suomi NPP/NOAA-20)",
         "Produtos utilizados:",
         "  - FIRMS: Deteccoes de focos ativos (pontos de calor)",
         "  - MCD64A1: Produto de area queimada mensal (500m)",
         "  - Sentinel-2: Imagens multiespectrais (10-20m)"])

    # SLIDE 12: Machine Learning
    add_content_slide(prs, "Aprendizado de Maquina: Gradient Boosting",
        ["Ensemble de arvores de decisao treinadas sequencialmente",
         "Cada arvore corrige erros das anteriores",
         "Algoritmos utilizados:",
         "  - LightGBM: Otimizado para velocidade e eficiencia",
         "  - XGBoost: Referencia em competicoes de ML",
         "Vantagens: Alta performance, interpretabilidade, robustez"])

    # SLIDE 13: Weak Labeling
    add_content_slide(prs, "Weak Labeling",
        ["Problema: Rotulacao manual e inviavel em larga escala",
         "Solucao: Gerar rotulos automaticamente usando dados auxiliares",
         "Abordagem proposta:",
         "  - Comparar deteccoes FIRMS com produto MCD64A1",
         "  - Janela temporal: +/- 15 dias",
         "  - Se ha cicatriz de queimada proxima: REAL",
         "  - Caso contrario: ESPURIA"])

    # SECAO 4: METODOLOGIA
    add_section_slide(prs, "Metodologia", 4)

    # SLIDE 15: Arquitetura
    add_content_slide(prs, "Arquitetura do Sistema",
        ["Pipeline em 4 etapas:",
         "1. INGESTAO: Download automatizado de dados brutos",
         "2. PROCESSAMENTO: Padronizacao e alinhamento espaciotemporal",
         "3. FEATURES: Extracao de caracteristicas + weak labeling",
         "4. MODELAGEM: Treinamento, validacao e inferencia"])

    # SLIDE 16: Fontes de Dados
    add_table_slide(prs, "Fontes de Dados Integradas",
        ["Fonte", "Tipo", "Resolucao", "Periodo"],
        [["FIRMS", "Focos ativos", "375m-1km", "2022-2024"],
         ["MCD64A1", "Area queimada", "500m", "Mensal"],
         ["Sentinel-2", "Multiespec.", "10-20m", "5 dias"],
         ["ERA5", "Meteorologia", "0.25 graus", "Horario"]])

    # SLIDE 17: Features Extraidas
    add_table_slide(prs, "Caracteristicas Extraidas",
        ["Categoria", "Features", "Descricao"],
        [["Temporal", "persistence_score, hour, month", "Padroes temporais"],
         ["Meteorologica", "temp, humidity, wind, precip", "Condicoes climaticas"],
         ["Vegetacao", "NDVI, NBR, dNBR", "Estado da vegetacao"],
         ["Espacial", "land_cover, distance_urban", "Contexto geografico"]])

    # SLIDE 18: Processo de Weak Labeling
    add_content_slide(prs, "Processo de Weak Labeling",
        ["Para cada deteccao FIRMS:",
         "1. Buscar pixels MCD64A1 em raio de 1km",
         "2. Verificar se data de queima esta em janela de +/-15 dias",
         "3. Atribuir rotulo:",
         "   - REAL (1): Correspondencia encontrada",
         "   - ESPURIA (0): Sem correspondencia",
         "Resultado: 708 deteccoes rotuladas automaticamente"])

    # SLIDE 19: Validacao
    add_content_slide(prs, "Estrategias de Validacao",
        ["Validacao Espacial (Leave-One-State-Out):",
         "  - Treina em 3 estados, testa no 4o",
         "  - Avalia generalizacao geografica",
         "",
         "Validacao Temporal:",
         "  - Treino: 2022-2023",
         "  - Teste: 2024",
         "  - Avalia estabilidade ao longo do tempo"])

    # SECAO 5: RESULTADOS
    add_section_slide(prs, "Resultados Experimentais", 5)

    # SLIDE 21: Metricas Principais
    add_metrics_slide(prs, "Desempenho do Classificador LightGBM",
        [("Acuracia", "82%"),
         ("ROC-AUC", "86%"),
         ("PR-AUC", "81%"),
         ("F1-Score", "80%")])

    # SLIDE 22: Comparacao Modelos
    add_table_slide(prs, "Comparacao: LightGBM vs XGBoost",
        ["Metrica", "LightGBM", "XGBoost"],
        [["Acuracia", "82%", "81%"],
         ["ROC-AUC", "86%", "85%"],
         ["PR-AUC", "81%", "80%"],
         ["Tempo Treino", "0.8s", "2.1s"]])

    # SLIDE 23: Validacao Espacial
    add_table_slide(prs, "Resultados da Validacao Espacial",
        ["Estado Teste", "Acuracia", "ROC-AUC", "PR-AUC"],
        [["Maranhao", "80%", "84%", "79%"],
         ["Tocantins", "83%", "87%", "82%"],
         ["Piaui", "81%", "85%", "80%"],
         ["Bahia", "82%", "86%", "81%"]])

    # SLIDE 24: Validacao Temporal
    add_content_slide(prs, "Resultados da Validacao Temporal",
        ["Treino (2022-2023) -> Teste (2024):",
         "",
         "  - Acuracia: 82% -> 77% (queda de 5 p.p.)",
         "  - ROC-AUC: 86% -> 81% (queda de 5 p.p.)",
         "",
         "Indica concept drift moderado",
         "Sugere necessidade de retreinamento periodico"])

    # SLIDE 25: Importancia Features
    add_content_slide(prs, "Importancia das Caracteristicas",
        ["Top 5 Features mais importantes:",
         "",
         "1. persistence_score (35%) - Persistencia temporal",
         "2. ndvi (18%) - Indice de vegetacao",
         "3. frp (12%) - Potencia radiativa do fogo",
         "4. temperature (9%) - Temperatura do ar",
         "5. brightness (7%) - Brilho do pixel"])

    # SLIDE 26: Impacto Operacional
    add_metrics_slide(prs, "Impacto Operacional Estimado",
        [("Reducao FP", "87%"),
         ("Alertas/Dia", "13/100"),
         ("Economia", "7x"),
         ("Recall", "95%")])

    # SECAO 6: CRONOGRAMA
    add_section_slide(prs, "Cronograma", 6)

    # SLIDE 28: Cronograma
    add_content_slide(prs, "Cronograma de Atividades",
        ["CONCLUIDO:",
         "  - Revisao bibliografica",
         "  - Implementacao do pipeline de dados",
         "  - Desenvolvimento da metodologia de weak labeling",
         "  - Treinamento e validacao dos modelos",
         "",
         "PENDENTE:",
         "  - Expansao do dataset",
         "  - Experimentos com deep learning",
         "  - Redacao final da dissertacao"])

    # SECAO 7: CONCLUSAO
    add_section_slide(prs, "Conclusoes", 7)

    # SLIDE 30: Contribuicoes
    add_content_slide(prs, "Contribuicoes do Trabalho",
        ["1. Pipeline de dados geoespaciais multifonte",
         "2. Metodologia de weak labeling para focos de incendio",
         "3. Conjunto de features para classificacao",
         "4. Modelo validado espacial e temporalmente",
         "5. Analise comparativa LightGBM vs XGBoost"])

    # SLIDE 31: Limitacoes
    add_content_slide(prs, "Limitacoes Identificadas",
        ["Dependencia do MCD64A1 (500m, mensal) para rotulacao",
         "Dataset relativamente pequeno (708 amostras)",
         "Especificidade regional (apenas MATOPIBA)",
         "Latencia dos dados ERA5",
         "Degradacao temporal requer retreinamento"])

    # SLIDE 32: Trabalhos Futuros
    add_content_slide(prs, "Trabalhos Futuros",
        ["Curto prazo:",
         "  - Expansao do dataset para mais anos",
         "  - Integracao com MapBiomas Fire",
         "",
         "Medio prazo:",
         "  - Modulo de predicao de propagacao (D+1)",
         "  - Arquiteturas de deep learning (CNN)",
         "",
         "Longo prazo:",
         "  - Sistema operacional de monitoramento",
         "  - Expansao para outras regioes (Amazonia, Pantanal)"])

    # SLIDE 33: Agradecimentos
    add_title_slide(prs,
        "Obrigado!",
        "Perguntas?\n\n[Seu Email]\n[Nome do Programa]")

    # Salvar apresentacao
    prs.save("Apresentacao_Qualificacao_Mestrado.pptx")
    print("Apresentacao salva: Apresentacao_Qualificacao_Mestrado.pptx")

def create_roteiro():
    """Cria o roteiro em DOCX"""
    doc = Document()

    # Titulo
    title = doc.add_heading("Roteiro para Apresentacao de Qualificacao", 0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER

    doc.add_paragraph("Classificacao Automatica de Deteccoes de Focos de Incendio na Regiao MATOPIBA Usando Aprendizado de Maquina")
    doc.add_paragraph("Tempo total estimado: 20-25 minutos")
    doc.add_paragraph("")

    # Roteiro por slide
    slides_roteiro = [
        ("Slide 1: Titulo", "30 seg", [
            "Cumprimentar a banca",
            "Apresentar-se brevemente",
            "Agradecer a presenca de todos"
        ]),
        ("Slide 2: Agenda", "30 seg", [
            "Apresentar a estrutura da apresentacao",
            "Mencionar que havera espaco para perguntas ao final"
        ]),
        ("Slide 3: Secao - Introducao", "10 seg", [
            "Transicao para o contexto do problema"
        ]),
        ("Slide 4: Contexto", "1.5 min", [
            "Brasil e um dos paises com maior incidencia de incendios",
            "INPE monitora via satelite e gera alertas automaticos",
            "Problema: muitos alertas sao falsos positivos",
            "Isso sobrecarrega os orgaos responsaveis",
            "Verificacao manual de cada alerta e impossivel"
        ]),
        ("Slide 5: Problema", "1 min", [
            "Explicar as fontes de falsos positivos",
            "Destacar o impacto: recursos desperdicados",
            "Motivar a necessidade de automacao"
        ]),
        ("Slide 6: Area de Estudo", "1 min", [
            "Apresentar a regiao MATOPIBA",
            "Explicar por que foi escolhida",
            "Mencionar disponibilidade de dados"
        ]),
        ("Slide 7: Secao - Objetivos", "10 seg", [
            "Transicao para objetivos"
        ]),
        ("Slide 8: Objetivo Geral", "30 seg", [
            "Ler o objetivo geral de forma clara",
            "Enfatizar: distinguir REAL de ESPURIA"
        ]),
        ("Slide 9: Objetivos Especificos", "1 min", [
            "Passar brevemente por cada objetivo",
            "Mencionar que todos foram alcancados"
        ]),
        ("Slide 10: Secao - Fundamentacao", "10 seg", [
            "Transicao para base teorica"
        ]),
        ("Slide 11: Sensoriamento Remoto", "1.5 min", [
            "Explicar o principio de deteccao termica",
            "Apresentar os satelites utilizados",
            "Descrever os produtos de dados"
        ]),
        ("Slide 12: Machine Learning", "1 min", [
            "Explicar gradient boosting de forma simples",
            "Destacar vantagens: performance e interpretabilidade",
            "Mencionar LightGBM como escolha principal"
        ]),
        ("Slide 13: Weak Labeling", "1 min", [
            "Explicar o problema de rotulacao manual",
            "Apresentar a solucao de weak labeling",
            "Detalhar a comparacao FIRMS vs MCD64A1"
        ]),
        ("Slide 14: Secao - Metodologia", "10 seg", [
            "Transicao para metodologia"
        ]),
        ("Slide 15: Arquitetura", "1 min", [
            "Apresentar as 4 etapas do pipeline",
            "Enfatizar a modularidade do sistema"
        ]),
        ("Slide 16: Fontes de Dados", "1 min", [
            "Descrever cada fonte brevemente",
            "Destacar a integracao multifonte"
        ]),
        ("Slide 17: Features", "1 min", [
            "Apresentar as categorias de features",
            "Mencionar total de 20+ caracteristicas"
        ]),
        ("Slide 18: Weak Labeling", "1 min", [
            "Detalhar o processo de rotulacao",
            "Mencionar resultado: 708 amostras"
        ]),
        ("Slide 19: Validacao", "1 min", [
            "Explicar validacao espacial (leave-one-state-out)",
            "Explicar validacao temporal",
            "Destacar o rigor metodologico"
        ]),
        ("Slide 20: Secao - Resultados", "10 seg", [
            "Transicao para resultados"
        ]),
        ("Slide 21: Metricas Principais", "1 min", [
            "Apresentar os 4 indicadores principais",
            "Destacar: 82% acuracia, 86% ROC-AUC",
            "Comparar com baseline (se houver)"
        ]),
        ("Slide 22: Comparacao Modelos", "30 seg", [
            "Mostrar que LightGBM e ligeiramente superior",
            "Destacar vantagem de tempo de treino"
        ]),
        ("Slide 23: Validacao Espacial", "1 min", [
            "Mostrar consistencia entre estados",
            "Variacao inferior a 4 pontos percentuais",
            "Indica boa generalizacao geografica"
        ]),
        ("Slide 24: Validacao Temporal", "1 min", [
            "Mostrar degradacao de ~5 p.p.",
            "Explicar concept drift",
            "Mencionar necessidade de retreinamento"
        ]),
        ("Slide 25: Importancia Features", "1 min", [
            "Destacar persistence_score como mais importante",
            "Explicar por que NDVI e relevante",
            "Insights para trabalhos futuros"
        ]),
        ("Slide 26: Impacto Operacional", "1 min", [
            "Destacar reducao de 87% em falsos positivos",
            "Traduzir em ganho pratico: 13 alertas vs 100",
            "Manter recall alto (95%)"
        ]),
        ("Slide 27: Secao - Cronograma", "10 seg", [
            "Transicao para cronograma"
        ]),
        ("Slide 28: Cronograma", "30 seg", [
            "Mostrar o que ja foi feito",
            "Apresentar proximos passos"
        ]),
        ("Slide 29: Secao - Conclusoes", "10 seg", [
            "Transicao para conclusoes"
        ]),
        ("Slide 30: Contribuicoes", "1 min", [
            "Enumerar as 5 contribuicoes principais",
            "Destacar originalidade do trabalho"
        ]),
        ("Slide 31: Limitacoes", "30 seg", [
            "Ser honesto sobre limitacoes",
            "Mostrar consciencia critica"
        ]),
        ("Slide 32: Trabalhos Futuros", "1 min", [
            "Apresentar direcoes futuras",
            "Mostrar visao de longo prazo"
        ]),
        ("Slide 33: Agradecimentos", "30 seg", [
            "Agradecer a banca novamente",
            "Abrir para perguntas"
        ])
    ]

    for slide_title, tempo, pontos in slides_roteiro:
        doc.add_heading(slide_title, level=1)
        doc.add_paragraph(f"Tempo: {tempo}")

        for ponto in pontos:
            p = doc.add_paragraph(ponto, style='List Bullet')

        doc.add_paragraph("")

    # Dicas finais
    doc.add_heading("Dicas para a Apresentacao", level=1)
    dicas = [
        "Pratique varias vezes cronometrando o tempo",
        "Mantenha contato visual com a banca",
        "Fale pausadamente e com clareza",
        "Esteja preparado para perguntas sobre metodologia e resultados",
        "Tenha os resultados numericos memorizados",
        "Leve agua para a apresentacao",
        "Chegue com antecedencia para testar o equipamento"
    ]
    for dica in dicas:
        doc.add_paragraph(dica, style='List Bullet')

    # Perguntas provaveis
    doc.add_heading("Perguntas Provaveis da Banca", level=1)
    perguntas = [
        "Por que escolheu LightGBM em vez de outros algoritmos?",
        "Como lidou com o desbalanceamento de classes?",
        "Qual a confiabilidade do weak labeling?",
        "Como o sistema se comportaria em outras regioes?",
        "Quais as limitacoes da resolucao do MCD64A1?",
        "Como pretende lidar com o concept drift em producao?",
        "Por que nao utilizou deep learning desde o inicio?",
        "Como validou que os rotulos do weak labeling estao corretos?"
    ]
    for pergunta in perguntas:
        doc.add_paragraph(pergunta, style='List Bullet')

    doc.save("Roteiro_Apresentacao_Qualificacao.docx")
    print("Roteiro salvo: Roteiro_Apresentacao_Qualificacao.docx")

if __name__ == "__main__":
    print("Criando apresentacao PPTX...")
    create_presentation()
    print("\nCriando roteiro DOCX...")
    create_roteiro()
    print("\nConcluido!")
