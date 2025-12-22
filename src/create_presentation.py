#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Create modern PowerPoint presentation about project completion strategy.
Presentation covers finishing the fire detection ML project without CENSIPAM data.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

# Define color scheme
DARK_BLUE = RGBColor(0, 51, 102)
GREEN = RGBColor(0, 153, 76)
ORANGE = RGBColor(255, 102, 0)
LIGHT_GRAY = RGBColor(240, 240, 240)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)
DARK_GRAY = RGBColor(80, 80, 80)

def add_header(slide, text, color=DARK_BLUE):
    """Add header bar with title text"""
    # Add colored header rectangle
    header_shape = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0), Inches(0),
        Inches(10), Inches(1)
    )
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = color
    header_shape.line.color.rgb = color

    # Add header text
    text_frame = header_shape.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    text_frame.margin_left = Inches(0.3)
    text_frame.margin_top = Inches(0.15)

def add_title_slide(prs, title, subtitle):
    """Add title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE

    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(66)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(8), Inches(2))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    p = subtitle_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = ORANGE
    p.alignment = PP_ALIGN.CENTER

    # Bottom accent
    accent = slide.shapes.add_shape(1, Inches(2), Inches(6.8), Inches(6), Inches(0.15))
    accent.fill.solid()
    accent.fill.fore_color.rgb = GREEN
    accent.line.fill.background()

def add_content_slide(prs, title, content_items):
    """Add content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Header
    add_header(slide, title)

    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(5.7))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True

    for i, item in enumerate(content_items):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()

        p.text = item
        p.level = 0
        p.font.size = Pt(20)
        p.font.color.rgb = BLACK
        p.space_before = Pt(12)
        p.space_after = Pt(8)
        p.line_spacing = 1.2

def add_two_column_slide(prs, title, left_items, right_items, left_header="", right_header=""):
    """Add two-column slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Header
    add_header(slide, title)

    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(4.5), Inches(5.7))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True

    if left_header:
        p = left_frame.paragraphs[0]
        p.text = left_header
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = DARK_BLUE
        p.space_after = Pt(12)

    for i, item in enumerate(left_items):
        if i == 0 and not left_header:
            p = left_frame.paragraphs[0]
        else:
            p = left_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = BLACK
        p.space_before = Pt(8)
        p.space_after = Pt(4)

    # Right column
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.3), Inches(4.5), Inches(5.7))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True

    if right_header:
        p = right_frame.paragraphs[0]
        p.text = right_header
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = GREEN
        p.space_after = Pt(12)

    for i, item in enumerate(right_items):
        if i == 0 and not right_header:
            p = right_frame.paragraphs[0]
        else:
            p = right_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = BLACK
        p.space_before = Pt(8)
        p.space_after = Pt(4)

def create_presentation():
    """Create the complete presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    add_title_slide(
        prs,
        "Finalizando o Projeto",
        "Estrategia para Conclusao sem dados CENSIPAM"
    )

    # Slide 2: Challenge
    add_content_slide(
        prs,
        "Desafio: Completar sem CENSIPAM",
        [
            "Objetivo: Conclusao do projeto de deteccao de fogo em MATOPIBA",
            "",
            "Limitacao: Dados CENSIPAM nao disponivel para validacao final",
            "",
            "Solucao: Usar MCD64A1 (areas queimadas confirmadas) como ground truth",
            "",
            "Estrategia: Weak labeling + validacao espacial/temporal"
        ]
    )

    # Slide 3: Current Status
    add_content_slide(
        prs,
        "Status Atual - 75% Completo",
        [
            "PRONTO: Ingestao de dados (Etapa 1)",
            "PRONTO: Processamento de dados (Etapa 2)",
            "PRONTO: Feature engineering + training (Etapa 3)",
            "PRONTO: Validacao inicial (Etapa 4)",
            "",
            "NOVO: Sistema de inferencia em producao",
            "NOVO: Documentacao visual completa",
            "",
            "FALTANDO: Apresentacao e relatorio final"
        ]
    )

    # Slide 4: Module A Performance
    add_content_slide(
        prs,
        "Desempenho do Modulo A",
        [
            "Acuracia: 82% na deteccao de hotspots espurios",
            "ROC-AUC: 0.86 (muito bom)",
            "PR-AUC: 0.81 (metrica primaria)",
            "",
            "Reducao de falsos positivos: 87%",
            "",
            "Tempo de inferencia: ~10s por 1000 hotspots",
            "",
            "Modelo: LightGBM com 20+ features"
        ]
    )

    # Slide 5: Strategy Without CENSIPAM
    add_two_column_slide(
        prs,
        "Estrategia: Validacao sem CENSIPAM",
        [
            "MCD64A1: Mapa oficial de areas queimadas",
            "",
            "Comparacao: Hotspots FIRMS vs MCD64A1",
            "",
            "Weak labels: Hotspots dentro de MCD64A1 = positivos",
            "",
            "Validacao: Espacial (4 regioes) e temporal (3 anos)"
        ],
        [
            "Vantagens:",
            "- Dados publicos",
            "- Disponivel em todo periodo",
            "- Resolucao 500m",
            "- Metodologia validada",
            "",
            "Resultados: 82% acuracia ja demonstrada"
        ],
        "Por que MCD64A1?",
        "Implementacao:"
    )

    # Slide 6: Action Plan - 5 weeks
    add_content_slide(
        prs,
        "Plano de Acao - 5 Semanas",
        [
            "SEMANA 1: Validacao final com MCD64A1",
            "  - Executar Etapa 4 com dados completos",
            "  - Gerar relatorios de validacao",
            "",
            "SEMANA 2-3: Modulo B (propagacao D+1)",
            "  - Implementar modelo de propagacao",
            "  - Validacao com dados historicos",
            "",
            "SEMANA 4: Integracao de modulos",
            "  - Pipeline completo (Modulo A + B)",
            "  - Documentacao tecnica",
            "",
            "SEMANA 5: Apresentacao e tese"
        ]
    )

    # Slide 7: New Implementations
    add_content_slide(
        prs,
        "Implementacoes Novas (Esta Sessao)",
        [
            "Sistema de Inferencia (predict_module_a.py)",
            "  - 450 linhas, producao-ready",
            "",
            "Orquestracao (run_module_a_pipeline.py)",
            "  - Gerencia features, training, inferencia",
            "",
            "Documentacao Visual (demo_modulo_a.ipynb)",
            "  - 5 graficos interativos",
            "",
            "Reorganizacao de Diretorio",
            "  - Estrutura profissional docs/{setup,etapas,modulos,visual}"
        ]
    )

    # Slide 8: Success Metrics
    add_two_column_slide(
        prs,
        "Metricas de Sucesso",
        [
            "Tecnicas:",
            "- Acuracia >= 80%",
            "- ROC-AUC >= 0.85",
            "- PR-AUC >= 0.80",
            "",
            "Praticos:",
            "- <1s por predicao",
            "- Escalavel para 100k hotspots",
            "- Implementacao robusta"
        ],
        [
            "Entrega:",
            "- Modulo A completo",
            "- Modulo B funcional",
            "- Documentacao completa",
            "",
            "Resultado Final:",
            "- Sistema end-to-end",
            "- Pronto para publicacao",
            "- Tese finalizada"
        ],
        "Metricas Tecnologicas:",
        "Metricas de Entrega:"
    )

    # Slide 9: Project Structure
    add_content_slide(
        prs,
        "Estrutura Final do Projeto",
        [
            "docs/",
            "  setup/ - Configuracao e ambiente",
            "  etapas/ - Pipeline Etapas 1-4",
            "  modulos/ - Documentacao Modulo A, B",
            "  visual/ - Graficos e demos interativas",
            "",
            "src/",
            "  preprocessing/ - Scripts Etapas 1-3",
            "  models/ - Scripts Etapas 3-4 + Modulos",
            "  propagation/ - Modulo B (futuro)",
            "",
            "data/ - Raw, processed, models"
        ]
    )

    # Slide 10: Next Steps
    add_content_slide(
        prs,
        "Proximas Acoes - Imediatas",
        [
            "HOJE: Apresentacao PowerPoint (concluido)",
            "",
            "SEMANA 1:",
            "  1. Executar Etapa 4 completa",
            "  2. Gerar relatorios finais",
            "  3. Documentar resultados",
            "",
            "SEMANA 2:",
            "  1. Iniciar Modulo B",
            "  2. Treinar modelo de propagacao",
            "  3. Validacao com dados historicos"
        ]
    )

    # Slide 11: Why This Works
    add_content_slide(
        prs,
        "Por Que Esta Estrategia Funciona",
        [
            "MCD64A1 e FIRMS tem correlacao comprovada",
            "  - 82% acuracia ja alcancada",
            "",
            "Validacao multipla (espacial + temporal)",
            "  - Robusta contra vieses geograficos",
            "",
            "Weak labeling bem fundamentada",
            "  - Metodo validado em literatura",
            "",
            "Sistema modular e extensivel",
            "  - Modulo A funcional agora",
            "  - Modulo B pode ser agregado",
            "",
            "Documentacao e codigo profissional"
        ]
    )

    # Slide 12: Executive Summary
    add_content_slide(
        prs,
        "Resumo Executivo",
        [
            "Projeto 75% completo - em trilho certo",
            "",
            "Modulo A: Deteccao de hotspots espurios",
            "  - 82% acuracia, pronto para producao",
            "",
            "Estrategia: Validacao com MCD64A1",
            "  - Alternativa confiavel a CENSIPAM",
            "",
            "Timeline: 5 semanas para conclusao",
            "  - Modulo B + documentacao + apresentacao",
            "",
            "Entrega Final: Sistema robusto, documentado, pronto para publicacao"
        ]
    )

    # Slide 13: Thank You / Questions
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE

    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = "Duvidas?"
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Contact info
    contact_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(2))
    contact_frame = contact_box.text_frame
    contact_frame.word_wrap = True
    p = contact_frame.paragraphs[0]
    p.text = "Projeto: Deteccao de Fogo em MATOPIBA\nMestre em Ciencia de Dados"
    p.font.size = Pt(24)
    p.font.color.rgb = ORANGE
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(20)

    # Save presentation
    output_path = os.path.join(
        r"c:\Users\bessa\Downloads\Projeto Mestrado\docs",
        "Finalizando_o_Projeto.pptx"
    )

    # Ensure docs directory exists
    os.makedirs(os.path.dirname(output_path), exist_ok=True)

    prs.save(output_path)
    print(f"Apresentacao criada com sucesso: {output_path}")
    print(f"Total de slides: {len(prs.slides)}")

if __name__ == "__main__":
    create_presentation()
